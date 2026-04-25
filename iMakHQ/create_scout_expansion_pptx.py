#!/usr/bin/env python3
"""メルカリスカウト 全カテゴリ展開計画 パワーポイント"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

DARK = RGBColor(0x1a, 0x1a, 0x2e)
BLUE = RGBColor(0x44, 0x72, 0xC4)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
YELLOW = RGBColor(0xF3, 0x9C, 0x12)
RED = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK2 = RGBColor(0x2C, 0x3E, 0x50)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
TEAL = RGBColor(0x1A, 0xBC, 0x9C)

def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK

def txt(slide, l, t, w, h, text, sz=18, c=WHITE, b=False, a=PP_ALIGN.LEFT):
    bx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b; p.alignment = a

def box(slide, l, t, w, h, fc, text="", sz=14, tc=WHITE):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fc; s.line.fill.background()
    if text:
        tf = s.text_frame; tf.word_wrap = True
        tf.paragraphs[0].text = text; tf.paragraphs[0].font.size = Pt(sz)
        tf.paragraphs[0].font.color.rgb = tc; tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def arrow(slide, l, t):
    txt(slide, l, t, 0.6, 0.6, "→", 28, YELLOW, True, PP_ALIGN.CENTER)

def down_arrow(slide, l, t):
    txt(slide, l, t, 0.6, 0.6, "↓", 28, YELLOW, True, PP_ALIGN.CENTER)

# ===== SLIDE 1: タイトル =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 2, 2, 12, 1.5, "メルカリスカウト", 48, WHITE, True, PP_ALIGN.CENTER)
txt(s, 2, 3.8, 12, 1, "全カテゴリ展開計画", 32, TEAL, True, PP_ALIGN.CENTER)
txt(s, 2, 5.5, 12, 1.2, "TCGで実証済みの仕入自動化を\nアパレル・G-SHOCK・一番くじに展開する", 20, GRAY, False, PP_ALIGN.CENTER)

# ===== SLIDE 2: 現在のTCGフロー（実証済み） =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.5, 0.3, 15, 0.8, "現在: TCG (PSA鑑定) — 実証済み", 28, GREEN, True)

box(s, 0.5, 1.5, 3, 1.5, PURPLE, "メルカリ検索\nPSA10 x シングル\nx 日本語版", 14)
arrow(s, 3.6, 1.9)
box(s, 4.2, 1.5, 3, 1.5, PURPLE, "画像解析\n(Claude API)\n「PSA番号を読め」", 14)
arrow(s, 7.3, 1.9)
box(s, 7.9, 1.5, 3, 1.5, BLUE, "eBay検索\nPSA番号で照合\n→ 確実にヒット", 14)
arrow(s, 11, 1.9)
box(s, 11.6, 1.5, 3, 1.5, GREEN, "GATE判定\nGO → 仕入候補\nNO-GO → スキップ", 14)

txt(s, 0.5, 3.5, 15, 0.5, "ポイント:", 18, WHITE, True)
points = [
    "PSA番号 = 世界でユニークなID → eBay検索が確実にヒットする",
    "画像にPSAラベルが写っている → Claude APIで番号を読み取れる",
    "PSA番号 → PSAサイト → 公式DB → Item Specifics全自動",
]
for i, p in enumerate(points):
    txt(s, 0.7, 4.1 + i * 0.45, 14, 0.4, f"・{p}", 14, GRAY)

# 実績
box(s, 0.5, 5.8, 15, 1.5, DARK2)
txt(s, 0.7, 5.9, 14, 0.5, "テスト実績:", 18, GREEN, True)
txt(s, 0.7, 6.4, 7, 0.4, "・15件巡回 → PSA番号読取 93%成功", 14, GRAY)
txt(s, 0.7, 6.8, 7, 0.4, "・GO 3件発見 (最大利益¥60,290)", 14, GRAY)
txt(s, 8, 6.4, 7, 0.4, "・画像解析コスト: ¥24/回", 14, GRAY)
txt(s, 8, 6.8, 7, 0.4, "・キャッシュで2回目以降コスト¥0", 14, GRAY)

# ===== SLIDE 3: 展開の課題 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.5, 0.3, 15, 0.8, "課題: なぜTCG以外に展開できていないか", 28, RED, True)

txt(s, 0.5, 1.3, 15, 0.5, "TCGが成功した理由:", 20, GREEN, True)
box(s, 0.5, 2, 7, 1.2, DARK2, "PSA番号 = ユニークID\n画像から確実に読み取れる\neBay検索で1発ヒット", 14)

txt(s, 0.5, 3.5, 15, 0.5, "他カテゴリの問題:", 20, RED, True)

cats = [
    ("アパレル\n(Porter/montbell)", "PSA番号がない\nタグに型番がないことも多い\n同一モデルでも色・サイズ違いで価格が変わる", "画像から何を読み取る？\n型番？ブランド名？商品名？"),
    ("G-SHOCK", "PSA番号がない\n型番 (GA-2100等) はタグに記載\nメルカリタイトルにも型番が入りやすい", "型番があればeBay検索は容易\n画像からの読取精度が課題"),
    ("一番くじ", "PSA番号がない\n景品名+キャラ名で特定\n「A賞」「ラストワン」等の用語", "景品名の日本語→英語変換が必要\n同一景品の状態差で価格変動大"),
]
for i, (cat, problem, question) in enumerate(cats):
    y = 4.2 + i * 1.5
    box(s, 0.5, y, 2.5, 1.2, ORANGE, cat, 12)
    txt(s, 3.2, y + 0.1, 5.5, 1, problem, 12, GRAY)
    txt(s, 9, y + 0.1, 6.5, 1, question, 12, YELLOW)

# ===== SLIDE 4: 解決策 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.5, 0.3, 15, 0.8, "解決策: Claude APIの画像解析を汎用化する", 28, WHITE, True)

# 共通アーキテクチャ
txt(s, 0.5, 1.2, 15, 0.5, "共通アーキテクチャ:", 20, TEAL, True)

box(s, 0.5, 1.9, 3, 1, PURPLE, "メルカリ商品画像", 14)
arrow(s, 3.6, 2.1)
box(s, 4.2, 1.9, 4.5, 1, PURPLE, "Claude API\n「この商品を特定して」", 14)
arrow(s, 8.8, 2.1)
box(s, 9.4, 1.9, 3, 1, BLUE, "eBay検索\n特定結果で照合", 14)
arrow(s, 12.5, 2.1)
box(s, 13.1, 1.9, 2.5, 1, GREEN, "GATE判定", 14)

# カテゴリ別のプロンプト
txt(s, 0.5, 3.3, 15, 0.5, "カテゴリ別: Claude APIに「何を読み取るか」を変えるだけ:", 18, WHITE, True)

prompts = [
    ("TCG (現在)", "「PSA認定番号を読み取れ」", "PSA #142490884", "PSA 10 142490884", GREEN),
    ("G-SHOCK", "「G-SHOCKの型番を読み取れ」", "GA-2100-1A1JF", "G-SHOCK GA-2100 PSA", YELLOW),
    ("アパレル", "「ブランド・商品名・型番を\n読み取れ」", "montbell Thunder Pass\n#1128344", "montbell Thunder Pass\nJacket", YELLOW),
    ("一番くじ", "「一番くじの景品名・キャラ名\nを読み取れ」", "ワンピース A賞\nルフィ フィギュア", "Ichiban Kuji One Piece\nLuffy Figure", YELLOW),
]
for i, (cat, prompt, result, ebay_query, color) in enumerate(prompts):
    y = 4 + i * 1.1
    box(s, 0.5, y, 2, 0.9, color, cat, 12)
    box(s, 2.7, y, 4, 0.9, DARK2, prompt, 11)
    arrow(s, 6.8, y + 0.15)
    box(s, 7.4, y, 3.5, 0.9, DARK2, result, 11)
    arrow(s, 11, y + 0.15)
    box(s, 11.6, y, 4, 0.9, DARK2, ebay_query, 11)

# ===== SLIDE 5: カテゴリ別実装計画 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.5, 0.3, 15, 0.8, "カテゴリ別 実装計画", 28, WHITE, True)

# G-SHOCK（最も簡単）
txt(s, 0.5, 1.2, 5, 0.5, "G-SHOCK — 難易度: 低", 20, GREEN, True)
items_gs = [
    "型番がタグ・裏蓋に必ず印字されている",
    "メルカリタイトルにも型番が含まれることが多い",
    "Claude API: 「型番を読み取れ」→ GA-2100-1A1JF",
    "eBay: 「G-SHOCK GA-2100」で検索 → ヒット率高い",
    "CASIO公式サイトでスペック確認可能",
]
for i, item in enumerate(items_gs):
    txt(s, 0.7, 1.7 + i * 0.38, 7, 0.35, f"・{item}", 12, GRAY)

# アパレル
txt(s, 8, 1.2, 7, 0.5, "アパレル — 難易度: 中", 20, YELLOW, True)
items_ap = [
    "型番がタグにないことがある（今朝のクレーム）",
    "ブランド名 + 商品名 + 色 + サイズで特定",
    "Claude API: 「ブランド・型番・サイズを読み取れ」",
    "eBay: 「montbell jacket [色] [サイズ]」で検索",
    "同一商品でも状態差で価格が変わる → 精度課題",
]
for i, item in enumerate(items_ap):
    txt(s, 8.2, 1.7 + i * 0.38, 7, 0.35, f"・{item}", 12, GRAY)

# 一番くじ
txt(s, 0.5, 3.8, 5, 0.5, "一番くじ — 難易度: 中高", 20, ORANGE, True)
items_ik = [
    "景品名 + キャラ名 + 賞 (A賞/ラストワン等) で特定",
    "日本語→英語の変換が必要（キャラ名・作品名）",
    "Claude API: 「景品名・キャラ・賞を読み取れ」",
    "eBay: 「Ichiban Kuji [作品] [キャラ] [賞]」で検索",
    "箱出し未開封 vs 開封品で価格差大 → 状態判定が課題",
]
for i, item in enumerate(items_ik):
    txt(s, 0.7, 4.3 + i * 0.38, 7, 0.35, f"・{item}", 12, GRAY)

# 優先順位
txt(s, 8, 3.8, 7, 0.5, "展開優先順位:", 20, WHITE, True)
priority = [
    ("1位", "G-SHOCK", "型番が確実に取れる。最も成功率が高い", GREEN),
    ("2位", "一番くじ", "景品名+キャラ名で特定可能。市場も大きい", YELLOW),
    ("3位", "アパレル", "型番なしの場合が多く、検索精度に課題", ORANGE),
]
for i, (rank, cat, reason, color) in enumerate(priority):
    y = 4.4 + i * 0.7
    box(s, 8, y, 1, 0.55, color, rank, 14)
    box(s, 9.1, y, 2, 0.55, DARK2, cat, 14)
    txt(s, 11.3, y + 0.05, 4.5, 0.5, reason, 12, GRAY)

# ===== SLIDE 6: 実装ステップ =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.5, 0.3, 15, 0.8, "実装ステップ: 何を変えるか", 28, WHITE, True)

txt(s, 0.5, 1.2, 15, 0.5, "mercari_scout.py の変更箇所:", 20, WHITE, True)

steps = [
    ("Step 1", "画像解析プロンプトの汎用化",
     "現在: 「PSA番号を読み取れ」固定\n変更: カテゴリに応じてプロンプトを切替\n  TCG → PSA番号 / G-SHOCK → 型番 / アパレル → ブランド+商品名",
     "小", GREEN),
    ("Step 2", "eBay検索クエリの汎用化",
     "現在: 「PSA 10 [cert番号]」固定\n変更: カテゴリに応じて検索クエリを生成\n  G-SHOCK → 「G-SHOCK [型番]」/ アパレル → 「[ブランド] [商品名]」",
     "小", GREEN),
    ("Step 3", "search_urls.txt にカテゴリタグ追加",
     "現在: URLだけ\n変更: 「# [category:gshock]」のようなタグを追加\n  → カテゴリに応じて画像解析とeBay検索を切替",
     "小", GREEN),
    ("Step 4", "GATE判定パラメータのカテゴリ対応",
     "現在: TCG固定 (FVF 13.25%, 送料¥2,000)\n変更: カテゴリごとに FVF・送料を切替\n  アパレル → FVF 15.3% / 一番くじ → 送料¥2,500",
     "小", GREEN),
]
for i, (step, title, desc, effort, color) in enumerate(steps):
    y = 1.9 + i * 1.5
    box(s, 0.5, y, 1.5, 1.2, color, step, 14)
    box(s, 2.1, y, 3, 1.2, DARK2, title, 13)
    txt(s, 5.3, y + 0.1, 9.5, 1.1, desc, 11, GRAY)
    box(s, 14.5, y + 0.3, 1.2, 0.5, color, f"工数:{effort}", 11)

txt(s, 0.5, 8, 15, 0.5, "全て「小」工数。既存のインフラ（eBay API、GATE判定、キャッシュ）はそのまま使える。", 16, YELLOW, True)

# ===== SLIDE 7: 完成後の姿 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.5, 0.3, 15, 0.8, "完成後の姿: 全カテゴリ自動仕入", 28, WHITE, True)

txt(s, 0.5, 1.2, 15, 0.8, "search_urls.txt に検索条件を1行追加するだけ", 24, TEAL, True, PP_ALIGN.CENTER)

# search_urls.txt の例
box(s, 1.5, 2.3, 13, 3.5, DARK2)
lines = [
    "# search_urls.txt",
    "",
    "# [category:tcg] ワンピースカード PSA10",
    "https://jp.mercari.com/search?category_id=1409&...",
    "",
    "# [category:tcg] ドラゴンボール PSA10",
    "https://jp.mercari.com/search?category_id=10861&...",
    "",
    "# [category:gshock] G-SHOCK",
    "https://jp.mercari.com/search?keyword=G-SHOCK&...",
    "",
    "# [category:apparel] montbell ジャケット",
    "https://jp.mercari.com/search?keyword=montbell&...",
    "",
    "# [category:ichiban] 一番くじ ワンピース",
    "https://jp.mercari.com/search?keyword=一番くじ+ワンピース&...",
]
for i, line in enumerate(lines):
    color = YELLOW if line.startswith("#") else GRAY if line else WHITE
    if line.startswith("https"):
        color = TEAL
    txt(s, 1.8, 2.4 + i * 0.2, 12.5, 0.2, line, 10, color)

txt(s, 0.5, 6.2, 15, 1, "python mercari_scout.py\n→ 全カテゴリを自動巡回 → 仕入GOリストを出力", 20, GREEN, True, PP_ALIGN.CENTER)

# ===== SLIDE 8: まとめ =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.5, 0.3, 15, 0.8, "まとめ", 28, WHITE, True)

summary = [
    ("実証済み", "TCG (PSA鑑定) でメルカリスカウトが稼働\n15件巡回 → GO 3件発見 → 最大利益¥60,290", GREEN),
    ("展開方針", "Claude APIの画像解析プロンプトを\nカテゴリごとに切り替えるだけ\n既存インフラはそのまま流用", BLUE),
    ("優先順位", "G-SHOCK (型番確実) → 一番くじ (景品名特定可)\n→ アパレル (型番なしの課題あり)", YELLOW),
    ("工数", "全Step「小」工数\nプロンプト変更 + 検索クエリ変更 + カテゴリタグ追加", GREEN),
    ("最終形", "search_urls.txtに1行追加 = 新カテゴリ参入\n人間は「何を売りたいか」だけ決める", TEAL),
]
for i, (title, desc, color) in enumerate(summary):
    y = 1.3 + i * 1.4
    box(s, 0.5, y, 2.5, 1.1, color, title, 18)
    txt(s, 3.3, y + 0.1, 12, 1, desc, 14, GRAY)

filepath = r"c:\Users\imax2\OneDrive\デスクトップ\iMak_Scout_Expansion.pptx"
prs.save(filepath)
print(f"保存: {filepath}")
