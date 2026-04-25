#!/usr/bin/env python3
"""iMak Trading Japan - eBay Listing System アーキテクチャ図 PPT 生成
レビュー用（Gemini想定）。スライド順:
1. 全体構成
2. データフロー
3. 4大論点と現状
4. 検証/品質保証レイヤー
5. ギャップ & ロードマップ
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

OUT = Path(__file__).parent / "system_architecture.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
BLUE = RGBColor(0x1F, 0x4E, 0x79)
GREEN = RGBColor(0x70, 0xAD, 0x47)
RED = RGBColor(0xC0, 0x50, 0x4D)
YELLOW = RGBColor(0xE8, 0xB4, 0x14)
GRAY = RGBColor(0x80, 0x80, 0x80)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xDE, 0xEB, 0xF7)
LIGHT_GREEN = RGBColor(0xE2, 0xEF, 0xDA)
LIGHT_RED = RGBColor(0xFB, 0xE5, 0xD6)


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BLUE
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.4))
    tfs = sub_box.text_frame
    tfs.text = subtitle
    p = tfs.paragraphs[0]
    p.font.size = Pt(13)
    p.font.color.rgb = GRAY
    return slide


def add_box(slide, x, y, w, h, text, fill_color=LIGHT_BLUE, font_size=10, bold=False, font_color=BLACK):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = BLUE
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.text = text
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = font_color
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=GRAY):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(1.5)
    return line


def add_label(slide, x, y, w, h, text, font_size=11, bold=True, color=BLUE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color


# ===========================================================
# Slide 1: 全体構成
# ===========================================================
s = add_title_slide(prs, "iMak Trading Japan - eBay Listing System 全体構成",
                    "目標: 月¥30万売上 / fb100 / 出品2000件 (現状: ¥1.5-5万 / fb50 / 842件)")

# データソース層
add_label(s, 0.3, 1.6, 3.0, 0.3, "① データソース", color=BLUE)
sources = [("Mercari", 0.3, 2.0), ("Mercari Shops", 0.3, 2.45), ("Rakuma", 0.3, 2.9),
           ("Amazon.co.jp", 0.3, 3.35), ("楽天市場", 0.3, 3.8), ("Yahoo!Shopping", 0.3, 4.25),
           ("UNIQLO公式API", 0.3, 4.7), ("CASIO公式", 0.3, 5.15), ("PSA cert", 0.3, 5.6)]
for name, x, y in sources:
    add_box(s, x, y, 1.7, 0.4, name, fill_color=LIGHT_BLUE, font_size=9)

# スカウト/取得層
add_label(s, 2.3, 1.6, 2.5, 0.3, "② スカウト", color=BLUE)
scouts = [("mercari_scout.py", 2.3, 2.0), ("montbell_outlet_\nscraper.py", 2.3, 2.6),
          ("手動URL貼付\n(Amazon/楽天等)", 2.3, 3.4)]
for name, x, y in scouts:
    add_box(s, x, y, 1.9, 0.55, name, fill_color=LIGHT_BLUE, font_size=9)

# スプシ層
add_label(s, 4.4, 1.6, 3.0, 0.3, "③ Google Sheets", color=BLUE)
sheets = [("統合High_商品管理\n(Tshirt/Porter/Kuji/Mont)", 4.4, 2.0),
          ("統合Low_商品管理\n(Reel/Tomica)", 4.4, 2.7),
          ("利益計算シートv2_GS\n(FVF/送料/為替)", 4.4, 3.4),
          ("処理共通化\nマトリクス.xlsx", 4.4, 4.1)]
for name, x, y in sheets:
    add_box(s, x, y, 2.5, 0.6, name, fill_color=LIGHT_GREEN, font_size=9)

# リスティング生成層
add_label(s, 7.1, 1.6, 3.0, 0.3, "④ リスティング生成", color=BLUE)
listers = [("tshirt_listing.py", 7.1, 2.0),
           ("mercari_to_ebay_csv.py\n--sheet porter/reel/tomica/kuji", 7.1, 2.5),
           ("montbell_listing.py", 7.1, 3.25),
           ("ichibankuji_to_csv.py", 7.1, 3.7),
           ("gshock_to_csv.py", 7.1, 4.15),
           ("psa_to_csv.py", 7.1, 4.6)]
for name, x, y in listers:
    add_box(s, x, y, 2.7, 0.45, name, fill_color=LIGHT_GREEN, font_size=9)

# 検証ライブラリ層
add_label(s, 9.95, 1.6, 3.0, 0.3, "⑤ 検証/共通ライブラリ", color=BLUE)
libs = [("whitelist_registry.py\n(7カテゴリenum/range)", 9.95, 2.0),
        ("listing_validator.py\n(構造+3AI議論)", 9.95, 2.7),
        ("listing_common.py\n(新設・未統合) ⚠", 9.95, 3.4),
        ("pricing_engine.py", 9.95, 4.1),
        ("profit_params.py (SSOT)", 9.95, 4.55)]
for name, x, y in libs:
    add_box(s, x, y, 3.05, 0.6, name, fill_color=LIGHT_BLUE, font_size=9)

# eBay 出品層
add_label(s, 0.3, 6.1, 3.0, 0.3, "⑥ eBay (US)", color=RED)
add_box(s, 0.3, 6.45, 3.0, 0.45, "File Exchange CSV upload", fill_color=LIGHT_RED, font_size=10)
add_box(s, 3.5, 6.45, 3.0, 0.45, "Trading API / Browse API", fill_color=LIGHT_RED, font_size=10)
add_box(s, 6.7, 6.45, 3.5, 0.45, "出品 → バイヤー視認", fill_color=LIGHT_RED, font_size=10, bold=True)
add_box(s, 10.4, 6.45, 2.6, 0.45, "在庫管理: トラホバx2", fill_color=LIGHT_GREEN, font_size=10)

# Arrows
add_arrow(s, 2.0, 3.5, 2.3, 3.5)
add_arrow(s, 4.2, 3.5, 4.4, 3.5)
add_arrow(s, 6.9, 3.5, 7.1, 3.5)
add_arrow(s, 9.8, 3.5, 9.95, 3.5)
add_arrow(s, 8.5, 5.1, 8.5, 6.45)

# ===========================================================
# Slide 2: 4大論点と現状
# ===========================================================
s = add_title_slide(prs, "4大論点と現状ステータス",
                    "インプレッション→売上 直結のコア4要素")

# 4 sections
sections = [
    ("① 価格設定", 0.3, "コスト+利益 (✓)\nFVF/送料 SSOT (✓)\n為替自動取得 (✓)\n為替ALERT (✓)\n\n❌ eBay競合median未取得\n❌ ALERT後HOLDなし\n❌ カテゴリ別最低価格なし",
     LIGHT_RED if True else LIGHT_GREEN, "60% 完了"),
    ("② タイトル", 3.55, "Claude API生成 (✓)\nSYSTEM_PROMPT指示 (✓)\nPython強制(reel) (✓)\nAmazon variation取得(✓)\n\n⚠ Title長保証 reelのみ\n⚠ Pre-owned/New整合 reelのみ\n⚠ 推測suffix防止 reelのみ\n⚠ 5スクリプト未対応",
     LIGHT_RED, "30% 完了"),
    ("③ Item Specifics", 6.8, "7カテゴリWhitelist (✓)\nリトライループ (✓)\nPlausibility range (✓ reel)\nmax_length (✓ reel)\nSource verification (✓ reel)\n\n⚠ 必須項目HOLD reel/porterのみ\n❌ eBay Required Fields照合無し\n⚠ 全カテゴリ展開未完",
     YELLOW, "50% 完了"),
    ("④ チェック+学習", 10.05, "ホワイトリスト検証 (✓)\nHOLDキュー (一部)\n3AI議論(PSA限定) (✓)\n\n❌ 回帰テスト無し\n❌ 過去指摘の自動再発防止無し\n❌ listing後audit自動化無し\n❌ iMakAudit活用ゼロ\n❌ AI review→ルール反映なし",
     LIGHT_RED, "20% 完了"),
]
for title, x, body, color, badge in sections:
    add_box(s, x, 1.6, 3.2, 0.5, title, fill_color=BLUE, font_size=14, bold=True, font_color=WHITE)
    add_box(s, x, 2.15, 3.2, 4.4, body, fill_color=color, font_size=10)
    add_box(s, x, 6.65, 3.2, 0.4, badge, fill_color=YELLOW, font_size=11, bold=True)

# ===========================================================
# Slide 3: 検証/品質保証フロー
# ===========================================================
s = add_title_slide(prs, "検証/品質保証フロー (現状 vs 目標)",
                    "Claude出力 → Python deterministic 強制 → CSV出力")

# 現状フロー (上段)
add_label(s, 0.3, 1.5, 3.0, 0.3, "【現状】", color=RED)
boxes_now = [
    ("Claude API\n生成", 0.3, 1.95),
    ("Whitelist\n検証\n+ retry", 1.85, 1.95),
    ("(reel限定)\nTitle整合\n強制", 3.4, 1.95),
    ("(reel限定)\nTitle\nパディング", 4.95, 1.95),
    ("HOLD\nキュー\n(部分)", 6.5, 1.95),
    ("CSV出力", 8.05, 1.95),
    ("eBay\nアップロード", 9.6, 1.95),
    ("人間目視\nレビュー", 11.15, 1.95),
]
for txt, x, y in boxes_now:
    color = LIGHT_RED if "(reel限定)" in txt or "部分" in txt else LIGHT_BLUE
    add_box(s, x, y, 1.45, 1.1, txt, fill_color=color, font_size=10)
for i in range(len(boxes_now) - 1):
    x1 = boxes_now[i][1] + 1.45
    x2 = boxes_now[i+1][1]
    add_arrow(s, x1, 2.5, x2, 2.5)

# 目標フロー (下段)
add_label(s, 0.3, 3.5, 3.0, 0.3, "【目標】", color=GREEN)
boxes_target = [
    ("Claude API\n生成", 0.3, 3.95),
    ("Whitelist\n検証\n+ retry\n(全カテゴリ)", 1.85, 3.95),
    ("listing_common\nTitle/SKU/\nコンディション\n強制", 3.4, 3.95),
    ("Plausibility\n+max_length\n+Source検証\n(全カテゴリ)", 4.95, 3.95),
    ("HOLD\nキュー\n(全カテゴリ)", 6.5, 3.95),
    ("回帰テスト\nfixtures.json\n物理ゲート", 8.05, 3.95),
    ("Final lint\naudit_csv_row()", 9.6, 3.95),
    ("CSV出力\n+ listing後\n自動audit", 11.15, 3.95),
]
for txt, x, y in boxes_target:
    add_box(s, x, y, 1.45, 1.4, txt, fill_color=LIGHT_GREEN, font_size=9)
for i in range(len(boxes_target) - 1):
    x1 = boxes_target[i][1] + 1.45
    x2 = boxes_target[i+1][1]
    add_arrow(s, x1, 4.65, x2, 4.65, color=GREEN)

# 学習レイヤー
add_label(s, 0.3, 5.6, 3.0, 0.3, "【学習機能 (新設)】", color=GREEN)
add_box(s, 0.3, 5.95, 6.2, 1.3,
        "improvement_log.jsonl\n指摘ログ蓄積\n→ 同指摘N回検知\n→ Pythonタスク化",
        fill_color=LIGHT_GREEN, font_size=10)
add_box(s, 6.7, 5.95, 6.3, 1.3,
        "fixtures_listing.json\n過去事例 (Baitcast Reel/PVC,ABS/1:65/9.7g/Pre-owned+New矛盾等)\n→ test_listing_rules.py で自動実行\n→ 失敗 = CSV出力ブロック",
        fill_color=LIGHT_GREEN, font_size=10)

# ===========================================================
# Slide 4: ギャップ & ロードマップ
# ===========================================================
s = add_title_slide(prs, "ギャップ & 実装ロードマップ",
                    "P0 = 即実装 / P1 = 1週間内 / P2 = 2週間以降")

# Header
header_y = 1.5
for x, w, h in [(0.3, 0.7, 0.4), (1.05, 7.5, 0.4), (8.6, 2.5, 0.4), (11.15, 1.85, 0.4)]:
    pass
add_box(s, 0.3, header_y, 0.7, 0.4, "優先", fill_color=BLUE, font_size=11, bold=True, font_color=WHITE)
add_box(s, 1.05, header_y, 7.5, 0.4, "実装内容", fill_color=BLUE, font_size=11, bold=True, font_color=WHITE)
add_box(s, 8.6, header_y, 2.5, 0.4, "影響カテゴリ", fill_color=BLUE, font_size=11, bold=True, font_color=WHITE)
add_box(s, 11.15, header_y, 1.85, 0.4, "期待効果", fill_color=BLUE, font_size=11, bold=True, font_color=WHITE)

rows = [
    ("P0", "listing_common.py 集約 (Title/SKU/コンディション/HOLD)", "全7カテゴリ", "水平展開漏れ解消"),
    ("P0", "5スクリプト書換 (tshirt/montbell/gshock/ichibankuji/...)", "全カテゴリ", "全listing統一品質"),
    ("P0", "回帰テスト fixtures + test_listing_rules.py", "全カテゴリ", "過去指摘再発0%"),
    ("P0", "audit_csv_row() final lint ゲート (CSV出力前)", "全カテゴリ", "違反物理ブロック"),
    ("P1", "eBay Browse API median取得 → pricing_engine 改修", "全カテゴリ", "価格適正化↑"),
    ("P1", "eBay Required Fields取得 → whitelist_registry反映", "全カテゴリ", "Failure防止"),
    ("P1", "improvement_log.jsonl + 同指摘N回検出", "全カテゴリ", "学習機能稼働"),
    ("P2", "control_panel から listing後 audit 自動実行", "全カテゴリ", "品質ゲート常時"),
    ("P2", "iMakAudit 監査官 listing毎自動呼出", "全カテゴリ", "既存仕組み活用"),
    ("P2", "AI review_logs → SYSTEM_PROMPT 自動反映", "全カテゴリ", "自動学習"),
]
y = header_y + 0.4
row_h = 0.42
for prio, content, cat, eff in rows:
    color = LIGHT_RED if prio == "P0" else LIGHT_BLUE if prio == "P1" else LIGHT_GREEN
    add_box(s, 0.3, y, 0.7, row_h, prio, fill_color=color, font_size=11, bold=True)
    add_box(s, 1.05, y, 7.5, row_h, content, fill_color=WHITE, font_size=10)
    add_box(s, 8.6, y, 2.5, row_h, cat, fill_color=WHITE, font_size=10)
    add_box(s, 11.15, y, 1.85, row_h, eff, fill_color=WHITE, font_size=10)
    y += row_h

# ===========================================================
# Slide 5: Gemini向けレビュー観点
# ===========================================================
s = add_title_slide(prs, "Gemini レビュー観点",
                    "以下の観点でフィードバック願います")

points = [
    "1. アーキテクチャ的な見落とし\n   (バッチ処理 vs リアルタイム、エラーリカバリ機構等)",
    "2. 品質保証の仕組みが目標 (月¥30万売上) に対して妥当か\n   現状ROI測定: 全カテゴリ整備で月+15-40%期待、4要素確実化で+50%目標",
    "3. listing_common.py 集約方針のリスク\n   5スクリプト書換 + 回帰テスト + 段階移行 で安全か",
    "4. 回帰テスト fixture設計案\n   カテゴリ別 fixture vs カテゴリ横断 fixture どちらが運用しやすいか",
    "5. eBay API活用の改善余地\n   Browse API median取得、Trading API Required Fields取得、Sold価格分析",
    "6. コスト面の懸念\n   Claude API リトライ × カテゴリ拡張で API 費用累積\n   (現状月$5-10、5x拡張で月$25-50想定)",
    "7. 見落としてる失敗モード\n   今日の作業で気づいてない品質リスク・運用リスク",
    "8. iMakAudit (Claude+Gemini二段監査官) の活用方針\n   既存仕組みを listing毎に自動呼出する vs 段階審査する 等",
]
y = 1.5
for p in points:
    add_box(s, 0.3, y, 12.7, 0.62, p, fill_color=LIGHT_BLUE, font_size=11)
    y += 0.7

prs.save(str(OUT))
print(f"OK 保存: {OUT}")
