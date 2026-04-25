#!/usr/bin/env python3
"""iMak Trading Japan システム構成図 パワーポイント生成 v2"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

DARK = RGBColor(0x1a, 0x1a, 0x2e)
BLUE = RGBColor(0x44, 0x72, 0xC4)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
YELLOW = RGBColor(0xF3, 0x9C, 0x12)
RED = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK2 = RGBColor(0x2C, 0x3E, 0x50)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
TEAL = RGBColor(0x1A, 0xBC, 0x9C)

def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK

def txt(slide, l, t, w, h, text, sz=18, c=WHITE, b=False, a=PP_ALIGN.LEFT):
    bx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b; p.alignment = a

def box(slide, l, t, w, h, fc, text="", sz=14, tc=WHITE):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fc; s.line.fill.background()
    if text:
        tf = s.text_frame; tf.word_wrap = True
        tf.paragraphs[0].text = text; tf.paragraphs[0].font.size = Pt(sz)
        tf.paragraphs[0].font.color.rgb = tc; tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def arrow(slide, l, t):
    txt(slide, l, t, 0.6, 0.6, "→", 28, YELLOW, True, PP_ALIGN.CENTER)

# ===== SLIDE 1: タイトル =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 2, 2, 12, 1.5, "iMak Trading Japan", 48, WHITE, True, PP_ALIGN.CENTER)
txt(s, 2, 3.8, 12, 1, "System Architecture & Roadmap", 28, GRAY, False, PP_ALIGN.CENTER)
txt(s, 2, 5, 12, 0.5, "2026-04-14", 20, GRAY, False, PP_ALIGN.CENTER)
txt(s, 2, 6.5, 12, 1, "仕入発見 → 選定 → 価格設定 → 品質チェック → 出品\nの全工程を自動化する仕組み", 18, TEAL, False, PP_ALIGN.CENTER)

# ===== SLIDE 2: カテゴリ全体像 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.5, 0.3, 15, 0.8, "取扱カテゴリ全体像", 28, WHITE, True)

cats = [
    ("TCG (PSA鑑定)", "One Piece / Pokemon\nDragon Ball / Gundam", "✅ 稼働中", GREEN, 0.5),
    ("アパレル", "montbell / Porter\nUNIQLO UT", "✅ 出品稼働\n🔧 スカウト未対応", ORANGE, 4),
    ("G-SHOCK", "CASIO G-SHOCK", "✅ 出品稼働\n🔧 スカウト未対応", ORANGE, 7.5),
    ("一番くじ", "フィギュア\n景品", "✅ 出品稼働\n🔧 スカウト未対応", ORANGE, 11),
]
for name, desc, status, color, x in cats:
    box(s, x, 1.3, 3.2, 1, color, name, 20)
    txt(s, x, 2.4, 3.2, 0.7, desc, 12, GRAY)
    txt(s, x, 3.2, 3.2, 0.7, status, 12, GREEN if "稼働中" in status else YELLOW)

# 共通基盤
txt(s, 0.5, 4.5, 15, 0.5, "全カテゴリ共通基盤:", 18, WHITE, True)
commons = [
    ("チェッカー (check_csv.py)", "タイトル・Item Specifics・送料バリデーション + AIレビュー", GREEN),
    ("GATE判定", "価格帯別の利益率・乖離率で GO/保留/NO-GO 自動判定", GREEN),
    ("eBay Browse API", "競合価格取得・TOPセラーItem Specifics参考表示", GREEN),
    ("market_log.csv", "全商品の市場データを蓄積（NO-GO含む）", GREEN),
]
for i, (name, desc, color) in enumerate(commons):
    y = 5.1 + i * 0.55
    box(s, 0.5, y, 3, 0.45, color, name, 12)
    txt(s, 3.7, y + 0.05, 12, 0.4, desc, 12, GRAY)

# ===== SLIDE 3: TCGフロー =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.5, 0.3, 15, 0.8, "TCG (PSA鑑定カード) フロー — 完全自動化済み", 28, WHITE, True)

# フロー
box(s, 0.3, 1.5, 2.8, 1.2, PURPLE, "メルカリスカウト\n検索条件巡回", 14)
txt(s, 0.3, 2.8, 2.8, 0.4, "PSA10絞り込み", 10, GRAY)
arrow(s, 3.2, 1.8)

box(s, 3.8, 1.5, 2.8, 1.2, PURPLE, "Claude API\nPSA番号読取", 14)
txt(s, 3.8, 2.8, 2.8, 0.4, "画像からcert番号特定", 10, GRAY)
arrow(s, 6.7, 1.8)

box(s, 7.3, 1.5, 2.8, 1.2, GREEN, "公式DB\nItem Specifics", 14)
txt(s, 7.3, 2.8, 2.8, 0.6, "Pokemon / OP / DB / Gundam\n4ゲーム全対応", 10, GRAY)
arrow(s, 10.2, 1.8)

box(s, 10.8, 1.5, 2.3, 1.2, BLUE, "GATE判定\n価格設定", 14)
txt(s, 10.8, 2.8, 2.3, 0.4, "GO/保留/NO-GO", 10, GRAY)
arrow(s, 13.2, 1.8)

box(s, 13.8, 1.5, 2, 1.2, ORANGE, "CSV出力\n+ チェック", 14)

# 公式DB詳細
txt(s, 0.3, 3.8, 15, 0.5, "公式データソース:", 18, WHITE, True)
dbs = [
    ("One Piece", "bandai_jp.py", "onepiece-cardgame.com", "Selenium"),
    ("Pokemon", "pokemon_card_jp.py", "pokemon-card.com", "Selenium"),
    ("Dragon Ball", "bandai_tcg_plus.py", "api.bandai-tcg-plus.com", "REST API"),
    ("Gundam", "bandai_tcg_plus.py", "api.bandai-tcg-plus.com", "REST API"),
]
for i, (game, module, url, method) in enumerate(dbs):
    y = 4.4 + i * 0.45
    box(s, 0.5, y, 2, 0.38, GREEN, game, 11)
    txt(s, 2.7, y, 3, 0.38, module, 11, YELLOW, True)
    txt(s, 5.8, y, 4, 0.38, url, 11, GRAY)
    txt(s, 10, y, 2, 0.38, method, 11, TEAL)

# ===== SLIDE 4: アパレル/G-SHOCK/一番くじ =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.5, 0.3, 15, 0.8, "アパレル / G-SHOCK / 一番くじ フロー", 28, WHITE, True)

# 現在
txt(s, 0.5, 1.3, 7, 0.5, "現在の仕組み (出品は稼働中):", 18, GREEN, True)

box(s, 0.5, 2, 3, 0.8, DARK2, "商品管理シート\n(手動入力)", 12)
arrow(s, 3.6, 2.1)
box(s, 4.2, 2, 3, 0.8, GREEN, "Claude API\n画像解析 → CSV", 12)
arrow(s, 7.3, 2.1)
box(s, 7.9, 2, 2.5, 0.8, BLUE, "GATE判定\n+ チェッカー", 12)
arrow(s, 10.5, 2.1)
box(s, 11.1, 2, 2, 0.8, ORANGE, "eBay出品", 12)

# 今後
txt(s, 0.5, 3.5, 7, 0.5, "今後の姿 (スカウト対応後):", 18, YELLOW, True)

box(s, 0.5, 4.2, 3, 0.8, PURPLE, "メルカリスカウト\n検索条件巡回", 12)
arrow(s, 3.6, 4.3)
box(s, 4.2, 4.2, 3, 0.8, PURPLE, "Claude API\n型番・商品名読取", 12)
arrow(s, 7.3, 4.3)
box(s, 7.9, 4.2, 2.5, 0.8, BLUE, "eBay照合\nGATE判定", 12)
arrow(s, 10.5, 4.3)
box(s, 11.1, 4.2, 2, 0.8, GREEN, "自動出品", 12)

# カテゴリ別の違い
txt(s, 0.5, 5.5, 15, 0.5, "カテゴリ別の画像解析:", 18, WHITE, True)
items = [
    ("TCG (今)", "画像 → PSA番号読取 → PSAサイトで特定", "✅ 稼働中"),
    ("アパレル", "画像 → ブランド+型番+サイズ読取 → eBay検索", "🔧 開発予定"),
    ("G-SHOCK", "画像 → 型番 (GA-2100等) 読取 → eBay検索", "🔧 開発予定"),
    ("一番くじ", "画像 → 景品名+キャラ名読取 → eBay検索", "🔧 開発予定"),
]
for i, (cat, desc, status) in enumerate(items):
    y = 6.1 + i * 0.45
    box(s, 0.5, y, 2, 0.38, TEAL, cat, 11)
    txt(s, 2.7, y, 9, 0.38, desc, 11, GRAY)
    txt(s, 12, y, 3, 0.38, status, 11, GREEN if "稼働" in status else YELLOW)

# ===== SLIDE 5: GATE判定 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.5, 0.3, 15, 0.8, "GATE判定: 価格帯別パラメータ (全カテゴリ共通)", 28, WHITE, True)

rows = [
    ["価格帯", "目標利益率", "許容乖離率", "意味"],
    ["$0-39", "25%", "50%", "低単価 → 率を上げて手間に見合う利益確保"],
    ["$40-100", "20-25%", "50%", "中低価格 → 積極出品で量を稼ぐ"],
    ["$100-200", "15%", "50%", "中価格 → 標準KPI"],
    ["$200-300", "10%", "40%", "中高価格 → 利益額は十分、率は緩和"],
    ["$300-500", "10%", "20-25%", "高価格 → 乖離を厳しく"],
    ["$500+", "10%", "10-15%", "超高額 → 最も厳格"],
]
for i, row in enumerate(rows):
    y = 1.3 + i * 0.55
    for j, cell in enumerate(row):
        x = 0.5 + j * 3.8
        fc = BLUE if i == 0 else DARK2
        box(s, x, y, 3.7, 0.45, fc, cell, 11 if j == 3 else 13)

txt(s, 0.5, 5.5, 15, 0.8, "出品判定:\n✅ GO (乖離≤0%) = 中央値x95%で出品  /  🟡 保留 (乖離≤許容%) = 目標価格で出品  /  ❌ NO-GO = CSV除外", 14, YELLOW)

# ===== SLIDE 6: メルカリスカウト詳細 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.5, 0.3, 15, 0.8, "メルカリスカウト: 仕入候補の自動発見", 28, WHITE, True)

# 実績
txt(s, 0.5, 1.2, 7, 0.5, "初回テスト実績 (ワンピース PSA10):", 18, GREEN, True)
results = [
    "15件巡回 → PSA番号読取 14/15成功 (93%)",
    "eBayヒット 4件 → GO 3件 + 保留 1件",
    "最大利益: ¥60,290 (トラファルガー・ロー)",
    "画像解析コスト: 約¥24/回 (15件)",
]
for i, r in enumerate(results):
    txt(s, 0.7, 1.7 + i * 0.4, 14, 0.4, f"・{r}", 14, GRAY)

# 機能一覧
txt(s, 0.5, 3.5, 7, 0.5, "搭載機能:", 18, WHITE, True)
features = [
    ("Chromeプロファイル", "メルカリログイン状態保持（パスキー対応）"),
    ("PSA番号キャッシュ", "同一商品の画像解析は1回だけ"),
    ("スプシ重複チェック", "出品済み商品を自動スキップ"),
    ("GOリストCSV", "go_list.csv → 仕入候補一覧"),
    ("certs_scout.txt", "psa_to_csv.pyに直接接続"),
    ("検索条件管理", "search_urls.txtにURL追加するだけ"),
]
for i, (name, desc) in enumerate(features):
    y = 4 + i * 0.42
    box(s, 0.5, y, 3, 0.35, TEAL, name, 11)
    txt(s, 3.7, y, 11, 0.35, desc, 11, GRAY)

# ===== SLIDE 7: 課題・ロードマップ =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.5, 0.3, 15, 0.8, "課題 & ロードマップ", 28, WHITE, True)

# 課題
txt(s, 0.5, 1.2, 7, 0.5, "現在の課題:", 20, RED, True)
issues = [
    ("メルカリスカウト", "TCG(PSA鑑定)のみ対応。アパレル/G-SHOCK/一番くじは未対応"),
    ("Dragon Ball/Gundam", "セットプレフィックス辞書が初期状態。新セット時に手動追加"),
    ("Pokemon辞書", "英名→和名辞書が限定的。未知のポケモンで公式DB検索失敗"),
    ("Finish/Attribute", "Claude API依存。判定精度の検証が必要"),
    ("サイズチャート", "衣類リスティングへの自動添付（GitHub画像アップロード待ち）"),
]
for i, (area, desc) in enumerate(issues):
    y = 1.8 + i * 0.45
    box(s, 0.5, y, 2.5, 0.38, RED, area, 11)
    txt(s, 3.2, y, 12, 0.38, desc, 11, GRAY)

# ロードマップ
txt(s, 0.5, 4.3, 7, 0.5, "ロードマップ:", 20, GREEN, True)
roadmap = [
    ("Phase 1 (完了)", "TCGリスティング全自動化 + メルカリスカウト(TCG)", GREEN),
    ("Phase 2 (次)", "メルカリスカウトを全カテゴリ対応\n画像→型番/商品名読取の汎用化", YELLOW),
    ("Phase 3", "Soldデータ連携（実売価格でGATE精度向上）", YELLOW),
    ("Phase 4", "バイヤー質問対応の自動化\n商品DB → 回答案自動生成", YELLOW),
    ("Phase 5", "全工程のワンコマンド化\nスカウト→仕入判断→出品→品質チェック", YELLOW),
]
for i, (phase, desc, color) in enumerate(roadmap):
    y = 4.9 + i * 0.65
    box(s, 0.5, y, 2.5, 0.55, color, phase, 12)
    txt(s, 3.2, y + 0.05, 12, 0.5, desc, 12, GRAY)

# ===== SLIDE 8: 定量効果 =====
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0.5, 0.3, 15, 0.8, "定量効果: Before → After", 28, WHITE, True)

items = [
    ("価格設定", "全商品$100均一", "市場中央値x95%\n価格帯別利益率", "同一5枚で¥17,867改善"),
    ("赤字チェック", "なし", "GATE判定で自動排除", "赤字出品ゼロ"),
    ("品質チェック", "Claudeに手動で聞く", "自動バリデーション\n+ AIレビュー", "作業時間 50分→0分"),
    ("仕入発見", "手動でメルカリ検索", "メルカリスカウト\nコマンド1発", "15件→GO 4件発見\n最大利益¥60,290"),
    ("Item Specifics", "手動入力・空欄多数", "公式DB自動取得\n4ゲーム対応", "根拠ある正確な値\nクレーム対策"),
]
for i, (item, before, after, effect) in enumerate(items):
    y = 1.3 + i * 1.3
    box(s, 0.5, y, 2.2, 0.9, DARK2, item, 13)
    box(s, 2.9, y, 3.3, 0.9, RED, before, 11)
    arrow(s, 6.3, y + 0.15)
    box(s, 7, y, 3.3, 0.9, GREEN, after, 11)
    box(s, 10.5, y, 5.2, 0.9, DARK2, effect, 11)

filepath = r"c:\Users\imax2\OneDrive\デスクトップ\iMak_System_Architecture.pptx"
prs.save(filepath)
print(f"保存: {filepath}")
