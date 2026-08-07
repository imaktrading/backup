#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""タイトル生成 & Item Specifics ロジック詳細 パワポ → デスクトップ。"""
import datetime
import os
import sys
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DESK = r"C:\Users\imax2\OneDrive\デスクトップ"
FONT = "Yu Gothic UI"
C_PRIMARY = RGBColor(0x1F, 0x4E, 0x79)
C_ACCENT = RGBColor(0xC0, 0x50, 0x00)
C_DARK = RGBColor(0x20, 0x28, 0x30)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GREY = RGBColor(0x80, 0x80, 0x80)
C_BG = RGBColor(0xF4, 0xF6, 0xF9)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = C_BG; bg.line.fill.background()
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def box(s, x, y, w, h, text, size=18, bold=False, color=C_DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    sp = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return sp


def header(s, title):
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_PRIMARY; bar.line.fill.background()
    box(s, 0.5, 0.1, 12.3, 0.72, title, size=24, bold=True, color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)


def table(s, x, y, w, rows, col_w, fs=12, hfs=12, row_h=0.42):
    gt = s.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w),
                            Inches(row_h * len(rows))).table
    for j, cw in enumerate(col_w):
        gt.columns[j].width = Inches(cw)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.margin_left = Inches(0.07); c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
            p = c.text_frame.paragraphs[0]; c.text_frame.word_wrap = True
            r = p.add_run(); r.text = str(val); r.font.name = FONT
            if i == 0:
                c.fill.solid(); c.fill.fore_color.rgb = C_PRIMARY
                r.font.size = Pt(hfs); r.font.bold = True; r.font.color.rgb = C_WHITE
            else:
                c.fill.solid(); c.fill.fore_color.rgb = C_WHITE if i % 2 else RGBColor(0xE8, 0xEE, 0xF5)
                r.font.size = Pt(fs); r.font.color.rgb = C_DARK
    return gt


# Slide 1 タイトル
s = slide()
hero = s.shapes.add_shape(1, 0, Inches(2.3), prs.slide_width, Inches(2.8))
hero.fill.solid(); hero.fill.fore_color.rgb = C_PRIMARY; hero.line.fill.background()
box(s, 0.8, 2.6, 11.7, 1.1, "タイトル生成 & Item Specifics ロジック詳細", size=34, bold=True, color=C_WHITE)
box(s, 0.8, 3.8, 11.7, 0.8, "出品くんが「売れる・正しい」タイトルと項目をどう作るか / eBay公式フィルタの活かし方",
    size=16, color=RGBColor(0xD9, 0xE2, 0xEC))
box(s, 0.8, 6.6, 11.7, 0.5, f"iMak Trading Japan  /  {datetime.date.today():%Y-%m-%d}", size=12, color=C_GREY)

# Slide 2 タイトル全体フロー
s = slide(); header(s, "① タイトル生成 — 全体フロー (4段)")
table(s, 0.5, 1.2, 12.3, [
    ["段", "関数", "やること"],
    ["1. 素地", "build_title()", "PSAの事実(キャラ+番号+set)を並べる。推論・改変なし"],
    ["2. 禁止語除去", "strip_banned_words()", "japanese/mint/graded/L@@K 等を除去 (Error 240回避)"],
    ["3. 字数調整", "pad_title()", "72〜80字に finish/card_type/set でキーワード補強"],
    ["4. SEO最適化", "refine_title()", "検索ヒット最大化の書換 (下記3フェーズ)。失敗時は元に戻す(耐障害)"],
], [2.0, 3.2, 7.1], row_h=0.6, fs=13)
box(s, 0.5, 4.4, 12.3, 1.0, "不変フィールド = character + card_number のみ (絶対に変えない)。\n"
    "それ以外 (set名・修飾語) は SEO 観点で書換可。タイトル先頭は常に \"PSA 10\"。",
    size=14, bold=True, color=C_ACCENT)

# Slide 3 build_title 詳細
s = slide(); header(s, "② タイトル素地 build_title — 事実ベース + SEO通称")
box(s, 0.5, 1.05, 12.3, 0.5, "構成式: 「PSA 10 + game_short + set_name + #card_number + subject」 → 80字超なら set名を落とす",
    size=14, bold=True, color=C_ACCENT)
box(s, 0.5, 1.65, 12.3, 0.45, "game_short = バイヤー検索通称 (SEO最大化)。一方 C:Game はeBay正規値 (dropdown hit) = 別軸戦略。", size=13)
table(s, 0.5, 2.2, 9.5, [
    ["公式名 (catalog)", "title用 game_short", "根拠"],
    ["Pokémon TCG", "Pokemon", "iMakKeywords Rank 1 (最強)"],
    ["Yu-Gi-Oh! TCG", "Yugioh", "Rank 19 (ハイフン/space無し最強)"],
    ["One Piece CCG", "One Piece", "Rank 13 (TCG suffix無しで節約)"],
    ["Dragon Ball SCG", "Dragon Ball SCG", "慣行略称"],
    ["Gundam Card Game", "Gundam TCG", "慣行略称"],
], [3.4, 3.0, 3.1], row_h=0.5, fs=12)
box(s, 0.5, 5.7, 12.3, 1.2, "例: 「PSA 10 Pokemon Crown Zenith #200 Mewtwo VSTAR」\n"
    "PSAのSubjectはsmart_titlecaseして使用。Mega系の重複(\"Mega…Mega Attack\")は分離して防止。",
    size=13, color=C_DARK)

# Slide 4 refine_title 3 phase
s = slide(); header(s, "③ タイトル SEO最適化 refine_title — 3フェーズ")
table(s, 0.5, 1.2, 12.3, [
    ["フェーズ", "内容", "データ源"],
    ["Phase 1", "eBay禁止語の置換 (Error 240回避) + technique→character 置換", "EBAY_FORBIDDEN_TERMS"],
    ["Phase 2", "上位キーワードでスコアリング (検索ボリューム加味して語を採用)", "iMakKeywords PDF (Q1 2026 ランキング)"],
    ["Phase 3", "TOPセラーのタイトル頻出語を抽出して反映", "sold_data xlsx (実売データ)"],
    ["最終", "複数 variants 生成 → 多角スコア → 最良案を採用", "(card_identification_agent と同形)"],
], [1.6, 6.4, 4.3], row_h=0.62, fs=12.5)
box(s, 0.5, 4.7, 12.3, 1.3, "守ること: 80字以内 / character + card_number は必ず残す / 重複語除去 / NGフィルタ。\n"
    "= 「売れるタイトル (検索ヒット)」と「eBay規約 (Error回避)」を両立。",
    size=13.5, color=C_DARK)

# Slide 5 Item Specifics 値の流れ
s = slide(); header(s, "④ Item Specifics — 値の流れ (catalog → 正規化 → fail-closed)")
table(s, 0.5, 1.15, 12.3, [
    ["段", "やること"],
    ["1. 値の源", "catalog(公式)の各フィールド → C:* 列にマッピング (Game/Set/Card Name/Character/Rarity/Card Type/Features/Language/Year/Country of Origin 等)"],
    ["2. 正規化", "whitelist_registry.validate_and_normalize() で eBay正規値へ (UNIQLO→Uniqlo, Tee→T-Shirt, Anime&Manga→Anime…)"],
    ["3. 自己修正", "違反は build_retry_feedback → Claude に再リクエスト (自己修正ループ)"],
    ["4. fail-closed", "確証なき値は空欄。Country of Origin は 'Does not apply' 明示 (eBayの勝手な補完=Japan を防ぐ)"],
], [1.5, 10.8], row_h=0.72, fs=12.5)
box(s, 0.5, 5.0, 12.3, 1.4, "カテゴリ別の必須Item Specifics (空欄=品質低下):\n"
    "  TCG: Game/Set/Card Name/Character/Rarity   G-SHOCK: Brand/Model/Type/Color/Movement\n"
    "  一番くじ: Brand/Character/Type/Franchise   Mercari(apparel): Brand/Type/Size/Color/Department",
    size=12.5, color=C_DARK)

# Slide 6 eBay公式フィルタ活用
s = slide(); header(s, "⑤ eBay公式フィルタ(Aspects)の活かし方 — 武器")
box(s, 0.5, 1.05, 12.3, 0.45, "eBayは各カテゴリの Item Specifics に『公式フィルタ値』を持つ。これを取得して検証に使う。",
    size=14, bold=True, color=C_ACCENT)
table(s, 0.5, 1.6, 12.3, [
    ["仕組み", "内容"],
    ["取得", "fetch_ebay_category_aspects.py (Sell Metadata API getItemAspectsForCategory) → ebay_<cat>_filter_lists_api.json (40 aspects)"],
    ["各aspectが持つ情報", "許容値リスト / aspect_usage(REQUIRED/RECOMMENDED/OPTIONAL) / aspect_mode(SELECTION_ONLY=リスト必須 / FREE_TEXT)"],
    ["監査で活用 (CSV監査くん)", "SELECTION_ONLY値が許容外→フィルタ不ヒット検出 / eBay推奨aspectの欠落→SEO機会"],
    ["生成で活用 (drift検査)", "手動whitelist が公式とズレてないか週次照合 (ズレ→正規化 or 修正)"],
], [3.2, 9.1], row_h=0.72, fs=12)
box(s, 0.5, 5.6, 12.3, 1.2, "効果: 「eBay自身が必須/推奨と言う項目」「正しい値か」を authoritative に判定 (TOPセラー推測より確実)。\n"
    "値は eBay正規値に合わせてこそフィルタにヒット = 検索で見つけてもらえる。",
    size=13, color=C_DARK)

# Slide 7 原則
s = slide(); header(s, "⑥ 設計原則まとめ")
box(s, 0.6, 1.3, 12.0, 4.5,
    "① 別軸戦略 — タイトル=バイヤー通称(検索最大化) / Item Specifics値=eBay正規値(フィルタhit)。役割を分ける。\n\n"
    "② 事実ベース + SEO — タイトルは PSA の事実を素地に、character+番号は不変、修飾語だけ SEO 書換。\n\n"
    "③ 値の捏造禁止 / fail-closed — catalog=正の辞書(SSOT)。確証なき値は空欄。誤った内容で出品しない。\n\n"
    "④ eBay公式フィルタを武器に — 取得した公式Aspectsで『正しい値か・必須/推奨か』を検証 (監査+生成drift)。\n\n"
    "⑤ 既存ロジックを再利用 — build_title/whitelist/check_csv を二重実装せず参照 (SSOT)。",
    size=15, color=C_DARK)

out = os.path.join(DESK if os.path.isdir(DESK) else ".",
                   f"タイトル_ItemSpecifics_ロジック詳細_{datetime.date.today():%Y%m%d}.pptx")
prs.save(out)
print("保存:", out)
