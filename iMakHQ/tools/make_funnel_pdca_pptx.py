#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""eBay Seller Hub レポート活用 PDCA サイクルの説明パワポを生成 → デスクトップ。

内容: ①4レポートで何が分かるか ②ファネル分析(症状→原因) ③役割分担 ④PDCA ⑤今月のアクション。
数値は listing_funnel の実機結果(2026-06-03 レポート)を反映。
"""
import os
import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DESK = r"C:\Users\imax2\OneDrive\デスクトップ"
FONT = "Meiryo"
C_BG = RGBColor(0xF4, 0xF6, 0xF8)
C_PRIMARY = RGBColor(0x1F, 0x4E, 0x79)
C_ACCENT = RGBColor(0xC0, 0x50, 0x2B)
C_GREEN = RGBColor(0x2E, 0x7D, 0x32)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK = RGBColor(0x22, 0x22, 0x22)
C_GREY = RGBColor(0x66, 0x66, 0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = C_BG; bg.line.fill.background()
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def box(s, x, y, w, h, text, size=18, bold=False, color=C_DARK, align=PP_ALIGN.LEFT,
        fill=None, anchor=MSO_ANCHOR.TOP):
    sp = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is not None:
        r = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        r.fill.solid(); r.fill.fore_color.rgb = fill; r.line.fill.background()
        s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(3, r._element)
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return sp


def header(s, title, num=None):
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_PRIMARY; bar.line.fill.background()
    box(s, 0.5, 0.12, 12.3, 0.76, title, size=26, bold=True, color=C_WHITE,
        anchor=MSO_ANCHOR.MIDDLE)


def table(s, x, y, w, rows, col_w, header_fill=C_PRIMARY, fs=13, hfs=13, row_h=0.42):
    n_rows = len(rows); n_cols = len(rows[0])
    gt = s.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(row_h * n_rows)).table
    for j, cw in enumerate(col_w):
        gt.columns[j].width = Inches(cw)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = gt.cell(i, j); c.margin_left = Inches(0.08); c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; r = p.add_run(); r.text = str(val)
            r.font.name = FONT
            if i == 0:
                c.fill.solid(); c.fill.fore_color.rgb = header_fill
                r.font.size = Pt(hfs); r.font.bold = True; r.font.color.rgb = C_WHITE
            else:
                c.fill.solid(); c.fill.fore_color.rgb = C_WHITE if i % 2 else RGBColor(0xEA, 0xEF, 0xF4)
                r.font.size = Pt(fs); r.font.color.rgb = C_DARK
    return gt


# ── Slide 1: タイトル ──
s = slide()
s.shapes.add_shape(1, 0, Inches(2.4), prs.slide_width, Inches(2.7)).fill.solid()
hero = s.shapes[-1]; hero.fill.fore_color.rgb = C_PRIMARY; hero.line.fill.background()
box(s, 0.8, 2.7, 11.7, 1.2, "eBay 出品改善 PDCA サイクル", size=40, bold=True, color=C_WHITE)
box(s, 0.8, 3.95, 11.7, 0.8, "Seller Hub 4レポート → ファネル分析 → 役割分担 → 改善ループ（API不要）",
    size=18, color=RGBColor(0xD9, 0xE2, 0xEC))
box(s, 0.8, 6.6, 11.7, 0.5, f"iMak Trading Japan  /  {datetime.date.today():%Y-%m-%d} 時点データ",
    size=12, color=C_GREY)

# ── Slide 2: 4レポートで何が分かるか ──
s = slide(); header(s, "① eBay レポート4種で何が分かるか")
box(s, 0.5, 1.1, 12.3, 0.5, "Seller Hub から DL する4つのレポートだけで、出品の全ファネルが見える（eBay API のクォータ不要）",
    size=15, bold=True, color=C_ACCENT)
table(s, 0.5, 1.75, 12.3, [
    ["レポート", "分かること", "キモの列"],
    ["全 active listings", "全4サイト(US/UK/AU/DE)の母集団・在庫数・watchers・過去販売数", "Available qty / Sold / Watchers / Site"],
    ["Listing quality report", "per-listing の【表示回数・CTR・転換率】+ 適正価格 + 項目欠落 + 写真/語数", "Impressions / CTR / Conversion / trending price"],
    ["unsold listings", "売れ残り・再出品(relist)状況", "Sold status / Relist status"],
    ["all orders", "実売（誰が・何を・いくらで買ったか）", "Item / Sale price / Buyer"],
], [2.6, 6.3, 3.4], row_h=0.78, fs=12)
box(s, 0.5, 5.9, 12.3, 1.2,
    "▶ 核心: Listing quality report が『表示→クリック→閲覧→購入』のファネルを per-listing で持つ。\n"
    "    これで「売れない」の症状を “検索に出てない / クリックされない / 買われない” の原因に分解できる。",
    size=15, bold=True, color=C_PRIMARY, fill=RGBColor(0xE8, 0xEF, 0xD8))

# ── Slide 3: ファネル分析 ──
s = slide(); header(s, "② ファネル分析：症状を「原因」に分解")
box(s, 0.5, 1.1, 12.3, 0.5, "在庫あり listing は3段階のどこで脱落したか / 在庫切れは需要で2分（今月の実数）",
    size=15, bold=True, color=C_ACCENT)
table(s, 0.5, 1.75, 12.3, [
    ["バケツ", "状態", "原因＝打つ手", "今月"],
    ["NO_SEARCH", "検索にほぼ出ない", "タイトルのキーワードが弱い", "84"],
    ["NO_CLICK", "表示有るがクリック0%", "サムネ/タイトル/価格", "114"],
    ["NO_CONVERT", "見られるが売れない", "価格/競合/説明", "280"],
    ["RESTOCK（在庫切れ）", "過去販売 or watcher 有", "需要実証済 → 再仕入れ", "237"],
    ["CULL（在庫切れ）", "一度も需要ゼロ", "出品停止を検討", "1,755"],
], [2.6, 3.2, 4.7, 1.8], row_h=0.62, fs=13)
box(s, 0.5, 6.05, 12.3, 1.0,
    "全 4,597 listing / 在庫切れ 1,992(43%) / 在庫あり 2,605 / US 深掘り対象 628。\n"
    "▶ 無在庫運用では『リピート可(G-SHOCK/Montbell)＝攻め』『1点もの(PSA/Porter)＝畳むか同カード再入手』で読み分ける。",
    size=13, color=C_DARK, fill=RGBColor(0xEF, 0xEF, 0xEF))

# ── Slide 4: 役割分担 ──
s = slide(); header(s, "③ 誰が・何を・どの技で（役割分担）")
table(s, 0.5, 1.2, 12.3, [
    ["バケツ", "対策", "担当システム", "技/ツール"],
    ["NO_SEARCH", "タイトル改修（検索語注入）", "Revise + キーワード", "iMakKeywords PDF 上位語"],
    ["NO_CLICK", "サムネ/画像の改善", "出品くん（HQ）", "実写サムネ / Vision 確認"],
    ["NO_CONVERT", "価格見直し → 値下げ/撤退", "抽出くん + HQ", "amazon_jp（原価）→ V8(US計算)"],
    ["RESTOCK", "仕入れ先を再確保", "抽出くん / 監視くん", "mercari_scout / Amazon在庫監視"],
    ["CULL", "出品停止 or 同カード再入手", "出品くん / 抽出くん", "end処理 / Mercari・スニダン照合"],
], [2.2, 3.6, 3.2, 3.3], row_h=0.66, fs=12)
box(s, 0.5, 5.9, 12.3, 1.1,
    "▶ 分析（出品くん=listing_funnel）が司令塔。各バケツを担当システムに割り振る。\n"
    "▶ 新規スクレイプは作らず “既存の技を借りる”（amazon_jp / mercari_scout / market_gate / Revise）。",
    size=14, bold=True, color=C_PRIMARY, fill=RGBColor(0xE8, 0xEF, 0xD8))

# ── Slide 5: PDCA ──
s = slide(); header(s, "④ PDCA サイクル（月次で回す）")
quad = [
    ("PLAN（計画）", C_PRIMARY, 0.5, 1.25,
     "月初に4レポートをDL\nlisting_funnel で5バケツに自動分類\n件数×単価で優先順位づけ"),
    ("DO（実行）", C_GREEN, 6.85, 1.25,
     "各担当が施策実行\n再仕入れ / タイトル / 価格 / 画像 / 停止\n（役割分担シートに沿って）"),
    ("CHECK（検証）", C_ACCENT, 0.5, 4.3,
     "翌月レポートを再DL\nimpr / CTR / 転換率 / 在庫切れ率 の変化を測定\nバケツ件数の増減で効果判定"),
    ("ACT（改善）", RGBColor(0x6A, 0x1B, 0x9A), 6.85, 4.3,
     "効いた施策を横展開\n効かない出品は撤退基準へ\nしきい値・タイトル型を更新"),
]
for title, col, x, y in [(q[0], q[1], q[2], q[3]) for q in quad]:
    pass
for title, col, x, y, body in quad:
    box(s, x, y, 5.95, 0.6, title, size=20, bold=True, color=C_WHITE, fill=col,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(s, x, y + 0.65, 5.95, 2.3, body, size=15, color=C_DARK, fill=C_WHITE)
box(s, 0.5, 7.0, 12.3, 0.4, "→ 同じレポートを毎月DLするだけで Check が自動で取れる＝改善が定量で回る",
    size=13, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)

# ── Slide 6: 今月のアクション ──
s = slide(); header(s, "⑤ 今月の優先アクション")
table(s, 0.5, 1.2, 12.3, [
    ["優先", "アクション", "対象", "効果/リスク"],
    ["1", "RESTOCK 再仕入れ（G-SHOCK35 / Montbell13）", "定番・需要実証済", "◎無在庫向き・即効・低リスク"],
    ["2", "NO_CONVERT 価格点検（Amazon原価→V8→値下げ/撤退）", "G-SHOCK 150", "売上直結・撤退判断も"],
    ["3", "CULL 整理（PSAは同カード再入手なら救出）", "1,755", "アカウント衛生・救出"],
    ["4", "NO_SEARCH タイトル改修（PDF→Revise）", "84", "中労力・検索流入"],
    ["5", "NO_CLICK 画像改善（高impr型番のサムネ）", "G-SHOCK 52", "クリック率改善"],
], [1.0, 6.2, 2.7, 2.4], row_h=0.62, fs=12)
box(s, 0.5, 5.95, 12.3, 1.0,
    "▶ 一推し: ①の再仕入れリスト（デスクトップ出力済）から着手。\n"
    "▶ 次回レポートで NO_CONVERT/在庫切れ件数が減れば施策が効いた証拠（=Check）。",
    size=14, bold=True, color=C_ACCENT, fill=RGBColor(0xEF, 0xEF, 0xEF))

out = os.path.join(DESK if os.path.isdir(DESK) else ".", f"eBayレポート活用_PDCAサイクル_{datetime.date.today():%Y%m%d}.pptx")
prs.save(out)
print("出力:", out, "/ slides:", len(prs.slides._sldIdLst))
