#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""在庫あり/在庫なし で分けた PDCA サイクルのパワポを生成 → デスクトップ。

在庫状態で対応が分かれる: 在庫あり=①②③(funnel診断)、在庫なし=④⑤(再仕入れ/畳む)。
2つの PDCA とライフサイクル(在庫なし→④→在庫あり→①②③)を図示。
"""
import datetime
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DESK = r"C:\Users\imax2\OneDrive\デスクトップ"
FONT = "Meiryo"
C_BG = RGBColor(0xF4, 0xF6, 0xF8)
C_PRIMARY = RGBColor(0x1F, 0x4E, 0x79)
C_INSTOCK = RGBColor(0x2E, 0x7D, 0x32)   # 在庫あり=緑
C_OOS = RGBColor(0xC0, 0x50, 0x2B)        # 在庫なし=橙
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK = RGBColor(0x22, 0x22, 0x22)
C_GREY = RGBColor(0x66, 0x66, 0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = C_BG; bg.line.fill.background()
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def box(s, x, y, w, h, text, size=18, bold=False, color=C_DARK, align=PP_ALIGN.LEFT,
        fill=None, anchor=MSO_ANCHOR.TOP, line=None):
    if fill is not None or line is not None:
        r = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        if fill is not None:
            r.fill.solid(); r.fill.fore_color.rgb = fill
        else:
            r.fill.background()
        if line is not None:
            r.line.color.rgb = line; r.line.width = Pt(1.5)
        else:
            r.line.fill.background()
        s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(3, r._element)
    sp = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run(); run.text = ln
        run.font.name = FONT; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    return sp


def header(s, title):
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_PRIMARY; bar.line.fill.background()
    box(s, 0.5, 0.12, 12.3, 0.76, title, size=25, bold=True, color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)


def table(s, x, y, w, rows, col_w, fs=12, hfs=12, row_h=0.5, hfill=C_PRIMARY):
    gt = s.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w),
                            Inches(row_h * len(rows))).table
    for j, cw in enumerate(col_w):
        gt.columns[j].width = Inches(cw)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = gt.cell(i, j); c.margin_left = Inches(0.08); c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
            p = c.text_frame.paragraphs[0]; run = p.add_run(); run.text = str(val); run.font.name = FONT
            if i == 0:
                c.fill.solid(); c.fill.fore_color.rgb = hfill
                run.font.size = Pt(hfs); run.font.bold = True; run.font.color.rgb = C_WHITE
            else:
                c.fill.solid(); c.fill.fore_color.rgb = C_WHITE if i % 2 else RGBColor(0xEA, 0xEF, 0xF4)
                run.font.size = Pt(fs); run.font.color.rgb = C_DARK
    return gt


def pdca_quad(s, x0, y0, w, color, title, items):
    """2x2 PDCA 図を (x0,y0) 起点 幅w で描画。items=[(label, body)]×4。"""
    box(s, x0, y0, w, 0.5, title, size=18, bold=True, color=C_WHITE, fill=color,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cw = (w - 0.2) / 2
    coords = [(x0, y0 + 0.6), (x0 + cw + 0.2, y0 + 0.6), (x0, y0 + 2.35), (x0 + cw + 0.2, y0 + 2.35)]
    for (cx, cy), (label, body) in zip(coords, items):
        box(s, cx, cy, cw, 0.4, label, size=13, bold=True, color=C_WHITE, fill=color,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        box(s, cx, cy + 0.42, cw, 1.3, body, size=10.5, color=C_DARK, fill=C_WHITE, line=color)


# ── Slide 1: タイトル ──
s = slide()
hero = s.shapes.add_shape(1, 0, Inches(2.4), prs.slide_width, Inches(2.6))
hero.fill.solid(); hero.fill.fore_color.rgb = C_PRIMARY; hero.line.fill.background()
box(s, 0.8, 2.65, 11.7, 1.2, "在庫別 PDCA サイクル", size=40, bold=True, color=C_WHITE)
box(s, 0.8, 3.9, 11.7, 0.8, "在庫あり（①②③ 診断で改善） / 在庫なし（④⑤ 再仕入れ・畳む）", size=18,
    color=RGBColor(0xD9, 0xE2, 0xEC))
box(s, 0.8, 6.6, 11.7, 0.5, f"iMak Trading Japan  /  {datetime.date.today():%Y-%m-%d}", size=12, color=C_GREY)

# ── Slide 2: なぜ在庫で分けるか ──
s = slide(); header(s, "なぜ在庫で分けるのか（対応が変わる）")
box(s, 0.5, 1.1, 12.3, 0.55,
    "在庫なし品は eBay が検索から外す → 表示/CTR/転換率データが存在しない → ①②③で診断できない",
    size=15, bold=True, color=C_OOS)
table(s, 0.5, 1.8, 12.3, [
    ["№", "やること", "在庫あり", "在庫なし", "理由"],
    ["①", "NO_SEARCH タイトル改修", "✓ 適用", "× 不可", "表示データが無い"],
    ["②", "NO_CLICK 画像改善", "✓ 適用", "× 不可", "CTR が測れない"],
    ["③", "NO_CONVERT 価格/仕入れ", "✓ 適用", "× 不可", "転換率が無い"],
    ["④", "RESTOCK 再仕入れ", "—", "✓ 適用", "在庫なし専用"],
    ["⑤", "CULL 整理 or 救出", "—", "✓ 適用", "在庫なし専用"],
], [0.7, 4.0, 2.2, 2.2, 3.2], row_h=0.62, fs=13)
box(s, 0.5, 5.95, 12.3, 1.0,
    "今月の数: 在庫あり 2,605件（①②③で診断） / 在庫なし 1,992件（④⑤で判定）。\n"
    "▶ 結論: 別物ではなく『順番』。在庫なしは まず ④⑤、戻したら ①②③ の診断対象になる。",
    size=14, bold=True, color=C_PRIMARY, fill=RGBColor(0xE8, 0xEF, 0xD8))

# ── Slide 3: 在庫あり PDCA ──
s = slide(); header(s, "在庫あり の PDCA（①②③ 診断で改善）")
box(s, 0.5, 1.05, 12.3, 0.45, "対象 2,605件。Listing quality report の表示→クリック→転換で脱落段階を特定し直す",
    size=13, bold=True, color=C_INSTOCK)
pdca_quad(s, 0.6, 1.6, 12.1, C_INSTOCK, "在庫あり PDCA", [
    ("PLAN（計画）", "Listing quality report で ①NO_SEARCH ②NO_CLICK ③NO_CONVERT に仕分け（impr/CTR/転換率）"),
    ("DO（実行）", "①タイトルを真の検索語先頭へ\n②画像/サムネ改善\n③値下げ or 代替仕入れ or 撤退"),
    ("CHECK（検証）", "翌月レポートで impr / CTR / 転換率 の改善を測定。各バケツ件数の増減で効果判定"),
    ("ACT（改善）", "効いたタイトル型/価格を横展開。改善しない出品は撤退基準へ"),
])
box(s, 0.6, 6.95, 12.1, 0.4, "▶ データが在るからこそ『どこで失敗してるか』を直せる（在庫ありの強み）",
    size=12, bold=True, color=C_INSTOCK, align=PP_ALIGN.CENTER)

# ── Slide 4: 在庫なし PDCA ──
s = slide(); header(s, "在庫なし の PDCA（④⑤ 再仕入れ・畳む）")
box(s, 0.5, 1.05, 12.3, 0.45, "対象 1,992件。表示データが無いので需要シグナル（過去販売/watcher）で仕分ける",
    size=13, bold=True, color=C_OOS)
pdca_quad(s, 0.6, 1.6, 12.1, C_OOS, "在庫なし PDCA", [
    ("PLAN（計画）", "需要シグナルで仕分け: ④RESTOCK（販売/watch 有=237）/ ⑤CULL（需要皆無=1,755）"),
    ("DO（実行）", "④同カード/型番を Mercari・Amazon で再仕入れ→在庫復活\n⑤畳む。※戻す前に①タイトル③価格を整える"),
    ("CHECK（検証）", "再出品後の販売有無・在庫切れ率。CULL 後のアカウント健全性（Google Shopping 等）"),
    ("ACT（改善）", "売れ筋は供給先を固定＋在庫監視。二度と入手不可は完全撤退"),
])
box(s, 0.6, 6.95, 12.1, 0.4, "▶ 在庫なしは『戻すか畳むか』の二択。戻す品は①③を前倒しで整えると復帰が強い",
    size=12, bold=True, color=C_OOS, align=PP_ALIGN.CENTER)

# ── Slide 5: ライフサイクル ──
s = slide(); header(s, "2つは繋がる：在庫ライフサイクル")
# 中央の循環図 (テキストボックス + 矢印代わりの記号)
box(s, 0.8, 1.5, 3.4, 1.1, "在庫なし\n（1,992件）", size=16, bold=True, color=C_WHITE, fill=C_OOS,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, 4.5, 1.5, 4.3, 1.1, "④ 再仕入れ可？\nMercari/Amazon 照合", size=14, bold=True, color=C_DARK,
    fill=C_WHITE, line=C_OOS, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, 9.1, 0.9, 3.5, 0.95, "YES → 在庫あり\n（①②③ 診断対象に）", size=13, bold=True, color=C_WHITE,
    fill=C_INSTOCK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, 9.1, 2.15, 3.5, 0.9, "NO → ⑤ 畳む", size=13, bold=True, color=C_WHITE, fill=C_GREY,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, 4.2, 0.62, 0.4, 0.4, "→", size=20, bold=True, color=C_OOS)
box(s, 8.8, 0.62, 0.4, 0.4, "→", size=20, bold=True, color=C_OOS)
# 在庫ありループ
box(s, 9.1, 3.4, 3.5, 1.0, "在庫あり\n① ② ③ で改善", size=15, bold=True, color=C_WHITE, fill=C_INSTOCK,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, 4.5, 3.4, 4.3, 1.0, "売れる → 一旦在庫切れ\n（無在庫: 売れてから仕入れ）", size=13, color=C_DARK,
    fill=C_WHITE, line=C_INSTOCK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, 0.8, 3.4, 3.4, 1.0, "→ 在庫なしへ戻る\n（ループ）", size=13, bold=True, color=C_OOS,
    fill=C_WHITE, line=C_OOS, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, 0.6, 4.9, 12.1, 1.9,
    "前倒しのコツ（重要）:\n"
    "  ・① タイトル / ③ 価格 は在庫なしのうちに準備できる（手元にデータが在る）→ 再仕入れ即、強い状態でスタート\n"
    "  ・② 画像 と impression/CTR は live（在庫あり）でないと取れない → 戻してから診断\n"
    "  ・無在庫モデルでは『売れる＝一旦在庫切れ』が正常。重要なのは 売れ筋の供給先を固定し再仕入れ100%を保つこと",
    size=13, color=C_DARK, fill=RGBColor(0xE8, 0xEF, 0xD8))

out = os.path.join(DESK if os.path.isdir(DESK) else ".", f"在庫別PDCAサイクル_{datetime.date.today():%Y%m%d}.pptx")
prs.save(out)
print("出力:", out, "/ slides:", len(prs.slides._sldIdLst))
