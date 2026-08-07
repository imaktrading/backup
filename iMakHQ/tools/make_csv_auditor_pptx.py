#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""CSV監査くん「何の出品に何のチェック→どう対応するか」詳細パワポ → デスクトップ。"""
import datetime
import os
import sys
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DESK = r"C:\Users\imax2\OneDrive\デスクトップ"
FONT = "Yu Gothic UI"
C_PRIMARY = RGBColor(0x1F, 0x4E, 0x79)
C_ACCENT = RGBColor(0xC0, 0x50, 0x00)
C_DARK = RGBColor(0x20, 0x28, 0x30)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GREY = RGBColor(0x80, 0x80, 0x80)
C_BG = RGBColor(0xF4, 0xF6, 0xF9)
C_OK = RGBColor(0x1E, 0x7A, 0x3C)
C_WARN = RGBColor(0xB8, 0x6A, 0x00)
C_NG = RGBColor(0xB0, 0x2A, 0x2A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = C_BG; bg.line.fill.background()
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def box(s, x, y, w, h, text, size=18, bold=False, color=C_DARK, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP):
    sp = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return sp


def header(s, title):
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_PRIMARY; bar.line.fill.background()
    box(s, 0.5, 0.1, 12.3, 0.72, title, size=24, bold=True, color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)


def table(s, x, y, w, rows, col_w, fs=12, hfs=12, row_h=0.40):
    n_rows, n_cols = len(rows), len(rows[0])
    gt = s.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w),
                            Inches(row_h * n_rows)).table
    for j, cw in enumerate(col_w):
        gt.columns[j].width = Inches(cw)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.margin_left = Inches(0.07); c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; r = p.add_run(); r.text = str(val); r.font.name = FONT
            if i == 0:
                c.fill.solid(); c.fill.fore_color.rgb = C_PRIMARY
                r.font.size = Pt(hfs); r.font.bold = True; r.font.color.rgb = C_WHITE
            else:
                c.fill.solid(); c.fill.fore_color.rgb = C_WHITE if i % 2 else RGBColor(0xE8, 0xEE, 0xF5)
                r.font.size = Pt(fs); r.font.color.rgb = C_DARK
    return gt


# ── Slide 1: タイトル ──
s = slide()
hero = s.shapes.add_shape(1, 0, Inches(2.3), prs.slide_width, Inches(2.8))
hero.fill.solid(); hero.fill.fore_color.rgb = C_PRIMARY; hero.line.fill.background()
box(s, 0.8, 2.6, 11.7, 1.1, "🔍 CSV監査くん — チェック内容と対応 一覧", size=36, bold=True, color=C_WHITE)
box(s, 0.8, 3.8, 11.7, 0.8, "出品くんが作った全カテゴリの出品CSVを、出品前に監査して自動対応する", size=17, color=RGBColor(0xD9, 0xE2, 0xEC))
box(s, 0.8, 6.6, 11.7, 0.5, f"iMak Trading Japan  /  {datetime.date.today():%Y-%m-%d}", size=12, color=C_GREY)

# ── Slide 2: 出品時チェック vs CSV監査くん 比較 ──
s = slide(); header(s, "出品時チェック vs CSV監査くん — 同じ深さ ＋ 修正/依頼/全カテゴリ")
box(s, 0.5, 1.05, 12.3, 0.45, "出品時チェック(check_csv)の深い検査を全部やった上に、修正・依頼・全カテゴリ対応を足したのがCSV監査くん。",
    size=14, bold=True, color=C_ACCENT)
table(s, 0.5, 1.65, 12.3, [
    ["検査 / 機能", "出品時チェック (check_csv)", "CSV監査くん"],
    ["機械ルール (タイトル/spec/cert/カテゴリ)", "✓", "✓"],
    ["市場ゲート (価格 GO/RELAX/HOLD/NO-GO ＋ 利益計算)", "✓", "✓"],
    ["TOPセラー Item Specifics 比較 (SEO)", "✓", "✓ ＋ 未対応spec全拾い(穴埋め)"],
    ["Claude AI 総合レビュー", "✓", "✓"],
    ["適用カテゴリ", "PSA TCG だけ自動", "全カテゴリ (TCG/G-shock/一番くじ/Mercari/他)"],
    ["送料ポリシー 自動修正", "✗ (警告のみ)", "✓ 機械修正"],
    ["データ誤り → カタログ修正依頼", "✗", "✓ 依頼書 自動生成"],
    ["生成バグ → プログラム修正依頼", "✗ (除外のみ)", "✓ 報告書 自動生成"],
    ["起動", "出品時(CSV生成)に自動", "ボタンで任意のCSVに後から"],
], [4.7, 3.6, 4.0], row_h=0.46, fs=11.5, hfs=11.5)

# ── Slide 3: 対応の4分岐 (凡例) ──
s = slide(); header(s, "対応は4分岐 — 「値の捏造」は絶対しない (SSOT/参照のみ)")
box(s, 0.5, 1.05, 12.3, 0.5, "不適正の種類によって対応が変わる。CSVに値を書き足すのは『機械的に一意に決まる送料』だけ。",
    size=14, bold=True, color=C_ACCENT)
table(s, 0.5, 1.7, 12.3, [
    ["対応", "発動条件", "動作"],
    ["🔧 機械修正", "送料ポリシーが価格帯と不一致 (決定論的)",
     "価格から正しい送料Profileを再計算してCSVを自動書換 (backup付・捏造でない)"],
    ["❌ 除外＋カタログ依頼", "データ誤り (set誤マップ / 必須spec空)",
     "該当行をCSVから物理除外(出品しない) ＋ Catalogへ修正依頼書を自動生成"],
    ["❌ 除外＋プログラム依頼", "生成プログラムの出力不正 (title/カテゴリ/日本語混入 等)",
     "該当行を除外 ＋ 生成スクリプト修正の報告書を自動生成 (CSVは捏造修正しない)"],
    ["💡 SEO報告のみ", "改善余地 (キーワード不足 / 推奨spec空 / 競合spec)",
     "CSVは一切触らず、改善メモとして報告のみ"],
], [2.4, 4.0, 5.9], row_h=0.95, fs=12.5)
box(s, 0.5, 6.7, 12.3, 0.6, "起動: 出品くん「🆕新規出品」パネルの『🔍出品前チェック』ボタン / コマンド: python csv_auditor.py (最新CSV自動)",
    size=12, color=C_GREY)

# ── Slide 3: カテゴリ別 チェック適用範囲 ──
s = slide(); header(s, "何の出品に何のチェックを入れているか (カテゴリ別)")
table(s, 0.4, 1.15, 12.5, [
    ["カテゴリ (出品)", "監査方式", "タイトル", "Item Specifics", "価格(市場)", "送料ポリシー", "カタログ整合"],
    ["PSA TCG", "フル4軸", "✓ 全項目", "✓ 必須5", "✓", "✓ 自動修正", "✓ set誤マップ検出"],
    ["G-SHOCK", "フル4軸", "✓ 全項目", "✓ 必須5", "✓", "✓ 自動修正", "－"],
    ["一番くじ", "フル4軸", "✓ 全項目", "✓ 必須4", "✓", "✓ 自動修正", "－"],
    ["Mercari (UT/montbell/porter)", "フル4軸*", "✓ 全項目", "△ 必須5(空は報告のみ)", "✓", "△ 報告のみ", "－"],
    ["リール/トミカ/Workman/その他", "汎用", "日本語混入・80字超 のみ", "－", "－", "－", "－"],
], [3.0, 1.3, 1.9, 2.5, 1.1, 1.6, 1.9], row_h=0.62, fs=12, hfs=11.5)
box(s, 0.4, 5.4, 12.5, 1.4,
    "✓全項目チェック = 80字以内/70字以上・禁止ワード・重複・先頭ルール。\n"
    "* Mercari注: apparel共通の必須spec(Size/Department等)がバッグ等に合わず全滅するため『空でも除外せず報告のみ』。\n"
    "  送料はcheck_csv内でTシャツ(UT)固定値のため誤適用回避で『自動修正せず報告のみ』。\n"
    "汎用 = 専用check_csvが無いカテゴリ。誤出品に直結する普遍項目(日本語混入・タイトル長)だけを最低限チェック。",
    size=12.5, color=C_DARK)

# ── Slide 4: チェック項目 × 検出 × 対応 (詳細) ──
s = slide(); header(s, "チェック項目 × 検出内容 × 対応 (詳細)")
table(s, 0.4, 1.1, 12.5, [
    ["チェック項目", "検出内容", "対応"],
    ["市場ゲート(価格)", "競合中央値との乖離で GO/RELAX/HOLD/NO-GO 判定", "NO-GO=❌除外 / HOLD=💡報告(再価格は出品側)"],
    ["Claude AI 総合レビュー", "バッチ全体の品質をAIが所見 (タイトル/SEO/整合)", "💡 報告 (レポートに記載)"],
    ["送料ポリシー", "価格帯と送料Profileが不一致", "🔧 機械修正 (価格から再計算してCSV書換)"],
    ["カタログ整合 (TCG)", "Set↔番号total / Set世代↔Year 不整合 (set誤マップ)", "❌ 除外 ＋ カタログ修正依頼"],
    ["必須Item Specifics", "必須項目が空欄", "❌ 除外 ＋ カタログ依頼 (Mercariは💡報告のみ)"],
    ["PSA鑑定番号 (TCG)", "番号が空 / 非数字", "❌ 除外 (fail-closed)"],
    ["価格", "数値でない", "❌ 除外 (fail-closed)"],
    ["タイトル長(上限)", "80字超", "❌ 除外 ＋ プログラム修正依頼"],
    ["タイトル先頭 (TCG)", "'PSA 10' で始まらない", "❌ 除外 ＋ プログラム依頼"],
    ["禁止ワード", "japanese/mint/graded/L@@K 等 (MercariはJapan許可)", "❌ 除外 ＋ プログラム依頼"],
    ["日本語混入", "タイトルに日本語文字", "❌ 除外 ＋ プログラム依頼"],
    ["カテゴリ/状態ID", "eBayカテゴリ・ConditionIDが規定外", "❌ 除外 ＋ プログラム依頼"],
    ["タイトル長(下限)", "70字未満 (キーワード不足)", "💡 SEO報告のみ (行は残す)"],
    ["推奨Item Specifics", "任意項目が空", "💡 SEO報告のみ"],
    ["競合TOPセラーspec", "自分が未対応の項目 (--with-market時)", "💡 SEO報告のみ"],
], [2.6, 5.6, 4.3], row_h=0.40, fs=11.5, hfs=12)

# ── Slide 5: 必須Item Specifics + 原則 ──
s = slide(); header(s, "カテゴリ別 必須Item Specifics ＆ 設計原則")
box(s, 0.5, 1.05, 12.3, 0.4, "必須Item Specifics (空欄だと出品品質低下 → 上表の対応へ)", size=14, bold=True, color=C_ACCENT)
table(s, 0.5, 1.5, 8.6, [
    ["カテゴリ", "必須Item Specifics"],
    ["PSA TCG", "Game / Set / Card Name / Character / Rarity"],
    ["G-SHOCK", "Brand / Model / Type / Color / Movement"],
    ["一番くじ", "Brand / Character / Type / Franchise"],
    ["Mercari (apparel)", "Brand / Type / Size / Color / Department  ※空は報告のみ"],
], [2.4, 6.2], row_h=0.52, fs=12)
box(s, 0.5, 4.5, 12.3, 2.6,
    "設計原則 (絶対遵守)\n"
    "  ① 値の捏造禁止 — カタログ=正の辞書(SSOT)。CSVに値を書き足すのは『送料(価格から一意)』だけ。\n"
    "  ② fail-closed — 確証なき/データ誤りの行は『除外』に倒す。間違った内容で出品しない方を常に優先。\n"
    "  ③ 根本へフィードバック — データ誤り→Catalog修正依頼、生成バグ→プログラム修正依頼を自動生成。\n"
    "  ④ 既存check_csvを再利用 — 各カテゴリの検査ロジックを二重実装せず参照 (SSOT)。\n"
    "出力: 修正後CSV / 除外レポート / カタログ依頼書(requests/) / プログラム報告書(review_logs/) / 監査レポート",
    size=13, color=C_DARK)

out = os.path.join(DESK if os.path.isdir(DESK) else ".",
                   f"CSV監査くん_チェックと対応_{datetime.date.today():%Y%m%d}.pptx")
prs.save(out)
print("✅ 保存:", out)
