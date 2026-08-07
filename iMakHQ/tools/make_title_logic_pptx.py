#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""タイトル生成ロジック説明書 + タイトルCSV監査 → PPTX (デスクトップ)。
構成: ① タイトル生成ロジック (全カテゴリ + 作成例) / ② タイトルのCSV監査 (CSV監査くん)。
後日 itemsp 版を同構成で作る前提。"""
import datetime
import os
import sys
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DESK = r"C:\Users\imax2\OneDrive\デスクトップ"
FONT = "Yu Gothic UI"
C_PRIMARY = RGBColor(0x1F, 0x4E, 0x79)
C_ACCENT = RGBColor(0xC0, 0x50, 0x00)
C_GREEN = RGBColor(0x2E, 0x7D, 0x32)
C_DARK = RGBColor(0x20, 0x28, 0x30)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GREY = RGBColor(0x80, 0x80, 0x80)
C_BG = RGBColor(0xF4, 0xF6, 0xF9)
C_CODE = RGBColor(0x0B, 0x3D, 0x2E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = C_BG; bg.line.fill.background()
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def box(s, x, y, w, h, text, size=18, bold=False, color=C_DARK, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, fill=None):
    sp = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    if fill is not None:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill; sp.line.color.rgb = C_GREY
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return sp


def code(s, x, y, w, h, text, size=14):
    sp = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor(0xEA, 0xF3, 0xEE); sp.line.color.rgb = C_GREEN
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15)
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = ln
        r.font.name = "Consolas"; r.font.size = Pt(size); r.font.color.rgb = C_CODE
    return sp


def header(s, title, sub=None):
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_PRIMARY; bar.line.fill.background()
    box(s, 0.5, 0.08, 12.3, 0.78, title, size=23, bold=True, color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        box(s, 0.5, 1.02, 12.3, 0.4, sub, size=13, color=C_ACCENT, bold=True)


def table(s, x, y, w, rows, col_w, fs=12, hfs=12, row_h=0.42):
    gt = s.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w),
                            Inches(row_h * len(rows))).table
    gt.first_row = False; gt.horz_banding = False
    for j, cw in enumerate(col_w):
        gt.columns[j].width = Inches(cw)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.margin_left = Inches(0.07); c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]; c.text_frame.word_wrap = True
            r = p.add_run(); r.text = str(val); r.font.name = FONT
            if i == 0:
                c.fill.solid(); c.fill.fore_color.rgb = C_PRIMARY
                r.font.size = Pt(hfs); r.font.bold = True; r.font.color.rgb = C_WHITE
            else:
                c.fill.solid(); c.fill.fore_color.rgb = C_WHITE if i % 2 else RGBColor(0xE8, 0xEE, 0xF5)
                r.font.size = Pt(fs); r.font.color.rgb = C_DARK
    return gt


# ============ 表紙 ============
s = slide()
hero = s.shapes.add_shape(1, 0, Inches(2.2), prs.slide_width, Inches(2.9))
hero.fill.solid(); hero.fill.fore_color.rgb = C_PRIMARY; hero.line.fill.background()
box(s, 0.8, 2.5, 11.7, 1.1, "タイトル生成ロジック 説明書", size=36, bold=True, color=C_WHITE)
box(s, 0.8, 3.7, 11.7, 0.9, "① 各カテゴリのタイトルの作り方 (作成例つき)  /  ② タイトルのCSV監査 (CSV監査くん)",
    size=16, color=RGBColor(0xD9, 0xE2, 0xEC))
box(s, 0.8, 6.6, 11.7, 0.5, f"iMak Trading Japan  /  {datetime.date.today():%Y-%m-%d}  /  ※itemsp版を同構成で別途作成",
    size=12, color=C_GREY)

# ============ Part① 区切り ============
s = slide()
band = s.shapes.add_shape(1, 0, Inches(3.0), prs.slide_width, Inches(1.5))
band.fill.solid(); band.fill.fore_color.rgb = C_ACCENT; band.line.fill.background()
box(s, 0.8, 3.25, 11.7, 1.0, "Part ①  タイトル生成ロジック", size=30, bold=True, color=C_WHITE)

# ============ 全体像: 3方式 ============
s = slide(); header(s, "全体像 — タイトル生成は3方式", "カテゴリごとに作り方が違う。これが基本構造")
table(s, 0.5, 1.6, 12.3, [
    ["方式", "中身", "対象カテゴリ"],
    ["A. ルール組立", "公式フィールドを決まった式で並べる。決定論的・AI不使用で安定", "G-SHOCK / montbell / Workman"],
    ["B. Claude AI 生成", "AIが買い手の検索語込みで生成。ただし商品別フォーマットをpromptで縛る", "一番くじ / Porter / Tomica / リール / UNIQLO"],
    ["C. ハイブリッド", "ルールで素地を作り → AIでSEO磨き (refine)", "PSA TCG"],
], [2.6, 6.6, 3.1], row_h=0.78, fs=13)
box(s, 0.5, 4.5, 12.3, 1.6,
    "■ 共通の鉄則 (全方式)\n"
    "・80字以内 (超過は優先度の低い語から削る / 単語の途中で切らない)\n"
    "・不変ファクト (型番・キャラ・番号) は崩さない  ・禁止語除去 (japanese/mint等 → Error 240回避)\n"
    "・SEO源は2つ: iMakKeywords PDF (検索ランキング) + TOPセラーの実売語 (sold_data)",
    size=13, color=C_DARK)

# ============ 下限パディング ============
s = slide(); header(s, "下限パディング — 70-80字を“実ファクト”で埋める", "短いタイトルは生成側で作り直す (監査WARN頼みにしない)")
box(s, 0.5, 1.15, 12.3, 1.05,
    "仕組み: タイトルが70字未満なら Item Specifics(=catalog公式値) を優先順に挿入して70-80字に。\n"
    "       80字を超える候補・既にタイトルに在る値はスキップ。材料が無ければ伸ばさない (=捏造禁止)。",
    size=14, bold=True, color=C_DARK)
code(s, 0.5, 2.3, 12.3, 0.95,
     'pad 材料の優先順 (pad_keys_priority):\n'
     'Color / Material / Size / Style  →  Water Resistance / Movement(時計)  →\n'
     'Theme / Franchise(フィギュア)  →  Type / Activity / Season(アパレル)  →  Year / Series', size=12)
box(s, 0.5, 3.45, 6.0, 0.4, "全カテゴリ対応 (2026-06-08):", size=13, bold=True, color=C_ACCENT)
box(s, 0.5, 3.85, 6.0, 1.6,
    "・G-SHOCK / 一番くじ / montbell / Workman を\n  実ファクト投入で有効化\n"
    "・Mercari(porter/tomica/reel) は元々有効\n・例: 61→74字, 49→74字, 66→79字", size=12)
box(s, 6.7, 3.45, 6.1, 0.4, "材料を増やす3経路 (誰が/どう):", size=13, bold=True, color=C_ACCENT)
table(s, 6.7, 3.85, 6.1, [
    ["状況", "誰が / どう"],
    ["①配線漏れ\n(catalogに値あるのに未投入)", "HQ/worktree がコードで投入"],
    ["②catalog空\n(spec自体が無い)", "Catalog/Harvest に依頼 (SSOT)"],
    ["③捏造", "禁止 (事実でない語は足さない)"],
], [2.7, 3.4], row_h=0.62, fs=11)

# ============ TCG ============
s = slide(); header(s, "PSA TCG — ハイブリッド (ルール素地 + AI/SEO磨き)")
code(s, 0.5, 1.15, 12.3, 0.75, 'PSA 10 {game略称} {セット} #{番号} {キャラ+rarity}   ← 80字超ならセット名を落とす', size=14)
table(s, 0.5, 2.1, 7.3, [
    ["公式名 (catalog)", "title用 略称", "根拠"],
    ["Pokémon TCG", "Pokemon", "検索Rank 1"],
    ["Yu-Gi-Oh! TCG", "Yugioh", "Rank 19"],
    ["One Piece CCG", "One Piece", "Rank 13"],
    ["Dragon Ball SCG", "Dragon Ball SCG", "慣行"],
], [3.0, 2.4, 1.9], row_h=0.46, fs=12)
box(s, 8.1, 2.0, 4.7, 0.4, "4段処理:", size=13, bold=True, color=C_ACCENT)
box(s, 8.1, 2.4, 4.7, 2.2,
    "1. build_title (事実を並べる)\n2. 禁止語除去\n3. 72-80字にパディング\n4. refine_title でSEO磨き\n   (失敗時は元に戻す=耐障害)",
    size=13)
box(s, 0.5, 4.5, 12.3, 0.5, "作成例:", size=13, bold=True, color=C_ACCENT)
code(s, 0.5, 4.95, 12.3, 1.4,
     'PSA 10 Pokemon Crown Zenith #200 Mewtwo VSTAR\n'
     'PSA 10 One Piece Egghead #ST29-004 Sanji Alternate Art Card Japanese\n'
     '不変 = キャラ名 + 番号 (絶対変えない) / セット名・修飾語は SEO で書換可', size=13)

# ============ G-SHOCK ============
s = slide(); header(s, "G-SHOCK — ルール組立 (キーワード最適化)")
code(s, 0.5, 1.2, 12.3, 0.7, 'CASIO G-Shock {型番} {特徴} Mens {Digital/Analog} Watch {色} New', size=14)
box(s, 0.5, 2.05, 12.3, 1.5,
    "・特徴 = Metal Covered / GPS / Bluetooth / Tough Solar (字数に余裕があれば優先度順に追加)\n"
    "・狙うキーワード = casio g shock / mens watches / watch men (eBay 2026Q1 PDF準拠)\n"
    "・80字超過時の削り順:  特徴キーワード → 色  の順に落とす", size=14)
box(s, 0.5, 3.7, 12.3, 0.5, "作成例:", size=13, bold=True, color=C_ACCENT)
code(s, 0.5, 4.15, 12.3, 1.3,
     'CASIO G-Shock GA-2100-1A1 Mens Digital Watch Black New\n'
     'CASIO G-Shock GAW-100B-1A2JF Mens Analog & Digital Watch Black Resin New', size=14)

# ============ 一番くじ ============
s = slide(); header(s, "一番くじ — Claude AI 生成 (フォーマット縛り)")
code(s, 0.5, 1.2, 12.3, 0.7, 'Ichiban Kuji [IP/Series] [Prize] [Character] [Figure Type] Bandai New', size=14)
box(s, 0.5, 2.05, 12.3, 1.3,
    "・AIが各賞のタイトルを生成。promptに上記フォーマット + TOPセラーの詳細サブタイトル例を入れて学習\n"
    "・生成後に normalize_title + 80字 trim\n"
    "・Series名はサブタイトルまで含める (TOPセラー慣習)", size=14)
box(s, 0.5, 3.6, 12.3, 0.5, "作成例:", size=13, bold=True, color=C_ACCENT)
code(s, 0.5, 4.05, 12.3, 1.5,
     'Ichiban Kuji One Piece A Prize Monkey D Luffy Masterlise Figure Bandai New\n'
     'Ichiban Kuji My Hero Academia A Prize Izuku Midoriya Masterlise Figure New\n'
     'Ichiban Kuji Dragon Ball B Prize Vegeta Figure New', size=13)

# ============ Mercari (porter/tomica/reel) ============
s = slide(); header(s, "Mercari系 — Claude AI 生成 (商品別フォーマット)", "商品ごとに別フォーマットをpromptで指定")
table(s, 0.5, 1.6, 12.3, [
    ["商品", "フォーマット", "作成例"],
    ["Porter", "YOSHIDA PORTER [Series] [Style] [Size] [Color] Used Japan",
     "YOSHIDA PORTER Tanker Shoulder Bag S Black Used Japan"],
    ["Tomica", "Tomica No.[X] [Make] [Model] [Color] [Scale] Vintage Japan",
     "Tomica No.47 Blue Nissan Gloria Van 1:64 Vintage Japan"],
    ["リール", "[Brand] [Model型番そのまま] ... Reel ... Japan",
     "Daiwa Zillion SV TW 1000 Baitcast Reel New Japan"],
], [1.6, 6.4, 4.3], row_h=0.82, fs=12)
box(s, 0.5, 4.6, 12.3, 1.4,
    "・Porter: Brand固定 \"Porter\" (3,514件の主戦場。HEAD PORTERタグ視認時のみ \"HEAD PORTER\")\n"
    "・Tomica: スケールは 1:64 に正規化 (1:65は使わない)\n"
    "・リール: 型番(モデル番号)をタイトル/タグ/公式から一字一句そのまま。Reel Sizeを必ず入れる",
    size=13, color=C_DARK)

# ============ montbell / Workman / UNIQLO ============
s = slide(); header(s, "montbell / Workman / UNIQLO")
table(s, 0.5, 1.25, 12.3, [
    ["商品", "方式", "フォーマット / 作成例"],
    ["montbell", "純テンプレ (AI不使用)",
     "montbell {英名} {色} US {USサイズ} (JP {JPサイズ}) {状態} Japan\n例: montbell Ultra Light Shell Jacket Yellow US L (JP XL) Pre-owned Japan"],
    ["Workman", "ルール組立",
     "{英名} {特徴} Workman {サブブランド} Japan Limited New\n超過時: 特徴→name末尾 の順に削る (途中切断なし)"],
    ["UNIQLO", "Claude AI 生成",
     "UNIQLO UT [Collab] [Character] T-Shirt [Color] US [Size] NWT Japan\n例: UNIQLO UT Doraemon T-Shirt Black US M (JP L) NWT Japan"],
], [1.7, 2.5, 8.1], row_h=1.05, fs=12)
box(s, 0.5, 5.7, 12.3, 1.1,
    "■ UNIQLO 重要: Brand は \"Uniqlo\" 固定。\"UNIQLO UT\" は eBayに無いブランド名 → フィルタ不ヒット = 売上機会喪失。\n"
    "  検索戦略: UNIQLOは検索上位外 → Theme(Anime/Music) + Character Family(コラボ/キャラ) で拾う。",
    size=13, bold=True, color=C_ACCENT)

# ============ Part② 区切り ============
s = slide()
band = s.shapes.add_shape(1, 0, Inches(3.0), prs.slide_width, Inches(1.5))
band.fill.solid(); band.fill.fore_color.rgb = C_GREEN; band.line.fill.background()
box(s, 0.8, 3.25, 11.7, 1.0, "Part ②  タイトルのCSV監査 (CSV監査くん)", size=28, bold=True, color=C_WHITE)

# ============ 監査の考え方 ============
s = slide(); header(s, "監査の考え方 — 「生成ロジック通りに作れているか」", "長さ/禁止語だけ見るのは表面。生成ロジックへの準拠を検証する")
box(s, 0.5, 1.6, 12.3, 1.3,
    "生成ロジックは catalog の事実から「タイトル」と「Item Specifics」の両方を作る。\n"
    "→ 両者は同じ事実を語るはず。食い違い = 生成のミス。これを2層で検証する。", size=15, bold=True, color=C_DARK)
table(s, 0.5, 3.0, 12.3, [
    ["層", "検査名", "見るもの"],
    ["① 整合", "title_spec_consistency", "spec の重要ファクトがタイトルに反映されているか"],
    ["② 形式", "title_format_checks", "カテゴリ別タイトルの「形」(先頭・必須語・Brand正規値) に従っているか"],
    ["③ SEO", "title_seo_findings", "PDF上位検索語を活かせてるか (同CSV内で相対的に弱い行を報告)"],
], [1.6, 4.3, 6.4], row_h=0.55, fs=13)
box(s, 0.5, 5.3, 12.3, 1.3,
    "■ 処置方針: タイトル不一致/形式逸脱/SEO弱は「報告のみ」(行は除外しない=誤除外回避) → プログラム/カタログ修正依頼へ回す。\n"
    "  ※ 長さ超過・禁止語・日本語混入・PSA10先頭欠落 などは従来通り別途検査 (誤出品直結は除外)。",
    size=13, color=C_ACCENT, bold=True)

# ============ 基本検査 (タイトル安全) ============
s = slide(); header(s, "基本検査 (タイトル安全) — 誤出品を物理的に止める層", "2層(整合/形式)の前提となる土台の検査")
table(s, 0.5, 1.6, 12.3, [
    ["検査", "条件 / 検出", "重大度", "処置"],
    ["文字数 超過", "80字より長い", "ERROR", "除外 + プログラム修正依頼"],
    ["文字数 不足", "70字未満 (キーワード不足の可能性)", "WARN", "報告のみ (除外しない)"],
    ["禁止ワード", "japanese / mint / graded / L@@K 等", "ERROR", "除外 + 修正依頼"],
    ["日本語混入", "ひらがな・カタカナ・漢字を検出", "ERROR", "除外 + 修正依頼"],
    ["PSA 10 先頭欠落", "TCGのみ: 先頭が \"PSA 10\" でない", "ERROR", "除外 + 修正依頼"],
    ["カテゴリ/ConditionID", "規定値でない", "ERROR", "除外 + 修正依頼"],
], [2.4, 5.0, 1.6, 3.3], row_h=0.54, fs=12)
box(s, 0.5, 5.55, 12.3, 1.3,
    "■ 文字数不足は「生成側で作り直す」を実装済 (2026-06-08): 全カテゴリで下限パディング有効化 → 短いまま出さない。\n"
    "  監査WARNは“まだ材料が足りず<70の行”を拾う最後の網 (= catalog充実 or 配線追加の宿題を示す)。\n"
    "  ※ ERROR系(除外) = 誤出品に直結。禁止語除去 strip_banned_words / 80字超のset名落とし も生成側で実装済。",
    size=13, color=C_ACCENT, bold=True)

# ============ 整合チェック ============
s = slide(); header(s, "① 整合チェック — タイトル ↔ Item Specifics")
box(s, 0.5, 1.15, 12.3, 0.45, "spec に値が在るのにタイトルに反映されていない → 生成ロジック逸脱の疑い", size=14, bold=True, color=C_ACCENT)
table(s, 0.5, 1.75, 12.3, [
    ["カテゴリ", "照合する列", "照合方式", "狙い"],
    ["G-SHOCK", "C:Model / C:Display / C:Band Color", "全体一致", "Modelにシリーズ名混入を検出"],
    ["TCG", "C:Character", "先頭語一致", "Character汚染(セット名混入)で誤検出しない"],
    ["一番くじ", "C:Character", "先頭語一致", "同上"],
    ["Mercari", "C:Color", "全体一致", "色の食い違い検出"],
], [2.0, 4.4, 2.3, 3.6], row_h=0.62, fs=12)
box(s, 0.5, 5.6, 12.3, 1.1,
    "照合方式を列ごとに分けるのが肝:\n"
    "・全体一致 = 値全体がタイトルに必要 (例 'G-SHOCK G-LIDE' が無ければ検出)\n"
    "・先頭語一致 = 先頭語だけ確認 (例 Character='Togekiss V ...汚染' でも 'Togekiss' が在ればOK)",
    size=13)

# ============ 形式準拠チェック ============
s = slide(); header(s, "② 形式準拠チェック — カテゴリ別タイトルの「形」")
table(s, 0.5, 1.2, 12.3, [
    ["カテゴリ/商品", "検査内容"],
    ["TCG", "番号 # を含む (PSA 10 先頭は別途検査済)"],
    ["G-SHOCK", "\"CASIO G-Shock\" で始まる + \"Watch\" を含む"],
    ["一番くじ", "\"Ichiban Kuji\" で始まる"],
    ["Porter", "\"PORTER\" + (Used / Pre-owned) を含む / Brand = Porter or HEAD PORTER"],
    ["Tomica", "\"Tomica\" を含む / Brand = Tomica"],
    ["UNIQLO", "T-Shirt / Tee を含む / Brand = Uniqlo (UNIQLO UT を検出)"],
    ["montbell", "\"montbell\" で始まる"],
    ["Workman", "\"Workman\" を含む"],
    ["リール", "\"Reel\" を含む"],
], [2.4, 9.9], row_h=0.52, fs=12, hfs=12)
box(s, 0.5, 6.45, 12.3, 0.5,
    "Mercariは商品混在 → C:Brand で自動判別して各商品の形を適用。", size=13, bold=True, color=C_ACCENT)

# ============ 実例 defect ============
s = slide(); header(s, "実例 — 監査が検出した本物の defect", "従来の監査では素通りしていた = ミスの原因")
table(s, 0.5, 1.6, 12.3, [
    ["検出", "実例", "影響", "処置"],
    ["G-SHOCK C:Model =\nシリーズ名 (型番でない)", "C:Model='G-SHOCK G-LIDE'\n←タイトルは型番 GBX-100NS-1JF", "Modelフィルタ\n不ヒット→露出減", "プログラム/\nカタログ修正依頼"],
    ["TCG C:Character\n汚染", "'Togekiss V Legendary Heartbeat'\n'Corviknight Vmax Vmax Climax'", "Characterフィルタ\n不ヒット", "同上"],
], [3.0, 5.0, 2.3, 2.0], row_h=1.15, fs=12)
box(s, 0.5, 5.4, 12.3, 1.2,
    "■ いずれもタイトル自体は正常だが Item Specifics が誤り → 整合チェックが食い違いとして検出。\n"
    "  根本(生成 or catalog)は別途修正案件として記録済。誤出品ではないが SEO 機会損失。",
    size=13, color=C_DARK)

# ============ ③ SEO監査 (PDF参照) ============
s = slide(); header(s, "③ タイトルSEO監査 — iMakKeywords PDF 参照", "「PDF上位検索語を活かせてるか」を監査する")
box(s, 0.5, 1.15, 12.3, 1.0,
    "仕組み: PDFを pdftotext で .txt 化 (C:/dev/iMak_data/keywords/) → 上位検索語でタイトルを採点 →\n"
    "       同じCSV内で SEO が相対的に弱い行を報告 (report-only)。csv_auditorは実行時pdftotextを呼ばず静的txt参照。",
    size=14, bold=True, color=C_DARK)
box(s, 0.5, 2.3, 12.3, 1.15,
    "■ なぜ「個別語を足せ」と言わないか (安全設計):\n"
    "   PDF上位語には他ブランド名 (rolex/omega 等) が混ざる → それを足せと言うと誤キーワード=捏造になる。\n"
    "   だから個別提案はせず、PDFプールでの相対採点で『弱い行』だけ炙り出す (人が見て判断)。",
    size=13, color=C_ACCENT, bold=True)
table(s, 0.5, 3.65, 12.3, [
    ["カテゴリ", "参照PDF (txt化済)"],
    ["TCG", "Toys_Hobbies_2026Q1"],
    ["G-SHOCK", "Jewelry_Watches_2026Q1"],
    ["一番くじ", "Collectibles_2026Q1"],
    ["Mercari (商品混在)", "C:Brandで出し分け: porter/uniqlo/montbell/workman→Clothing / tomica→Toys / リール→Sporting"],
], [2.6, 9.7], row_h=0.5, fs=12)
box(s, 0.5, 6.35, 12.3, 0.7,
    "検証: G-SHOCK 全行均一(テンプレ生成)で誤フラグ0 / TCG One Piece も閾値内。Mercari は商品グループ毎に相対比較。",
    size=12, color=C_GREY)

# ============ まとめ ============
s = slide(); header(s, "まとめ")
box(s, 0.6, 1.4, 12.1, 4.8,
    "① タイトル生成ロジック\n"
    "   ・3方式: ルール組立(G-SHOCK/montbell/Workman) / AI生成(一番くじ/Mercari/UNIQLO) / ハイブリッド(TCG)\n"
    "   ・共通鉄則: 80字 / 不変ファクト保持 / 禁止語除去 / 下限70字パディング / SEO源=iMakKeywords PDF + TOPセラー\n\n"
    "② タイトルのCSV監査 (CSV監査くん) — 基本検査 + 3層:\n"
    "   ・基本: 80字超/70字未満/禁止語/日本語/PSA10先頭  ①整合(title↔Item Sp) ②形式(カテゴリ別の形) ③SEO(PDF上位語)\n"
    "   ・不一致/逸脱/SEO弱は報告→修正依頼 (誤除外しない)。実データで本物の defect を検出済\n\n"
    "→ 次: 同じ構成で『Item Specifics (itemsp) 版』を作成する",
    size=15, color=C_DARK)

out = os.path.join(DESK, f"タイトル生成ロジック説明書_{datetime.date.today():%Y%m%d}.pptx")
prs.save(out)
print(f"✅ 保存: {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
