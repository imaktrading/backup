#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""出品3状態 × 効果的アクション のパワポ → デスクトップ。

① 出品済・在庫あり=直す(効果小〜中) / ② 出品済・在庫なし=戻す/畳む(中) / ③ 新規出品=増やす(大)。
"""
import datetime
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DESK = r"C:\Users\imax2\OneDrive\デスクトップ"
FONT = "Meiryo"
C_BG = RGBColor(0xF4, 0xF6, 0xF8)
C_PRIMARY = RGBColor(0x1F, 0x4E, 0x79)
C_S1 = RGBColor(0x2E, 0x7D, 0x32)   # 在庫あり=緑
C_S2 = RGBColor(0xC0, 0x50, 0x2B)   # 在庫なし=橙
C_S3 = RGBColor(0x1F, 0x4E, 0x79)   # 新規=青(本丸)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK = RGBColor(0x22, 0x22, 0x22)
C_GREY = RGBColor(0x66, 0x66, 0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = C_BG; bg.line.fill.background()
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def box(s, x, y, w, h, text, size=18, bold=False, color=C_DARK, align=PP_ALIGN.LEFT,
        fill=None, anchor=MSO_ANCHOR.TOP, line=None):
    if fill is not None or line is not None:
        r = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        if fill is not None:
            r.fill.solid(); r.fill.fore_color.rgb = fill
        else:
            r.fill.background()
        if line is not None:
            r.line.color.rgb = line; r.line.width = Pt(1.5)
        else:
            r.line.fill.background()
        s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(3, r._element)
    sp = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(2)
        run = p.add_run(); run.text = ln
        run.font.name = FONT; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    return sp


def header(s, title):
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_PRIMARY; bar.line.fill.background()
    box(s, 0.5, 0.12, 12.3, 0.76, title, size=25, bold=True, color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)


def column(s, x, color, head, count, effect, actions, tools):
    w = 4.0
    box(s, x, 1.2, w, 0.95, head, size=17, bold=True, color=C_WHITE, fill=color,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(s, x, 2.2, w, 0.5, f"{count}    効果: {effect}", size=13, bold=True, color=color,
        align=PP_ALIGN.CENTER, fill=C_WHITE, line=color)
    box(s, x, 2.78, w, 3.3, actions, size=12.5, color=C_DARK, fill=C_WHITE, line=color)
    box(s, x, 6.15, w, 0.85, "ツール: " + tools, size=10.5, color=C_GREY, fill=RGBColor(0xEE, 0xEE, 0xEE))


# ── Slide 1: タイトル ──
s = slide()
hero = s.shapes.add_shape(1, 0, Inches(2.4), prs.slide_width, Inches(2.6))
hero.fill.solid(); hero.fill.fore_color.rgb = C_PRIMARY; hero.line.fill.background()
box(s, 0.8, 2.65, 11.7, 1.2, "出品 3状態 × 効果的アクション", size=38, bold=True, color=C_WHITE)
box(s, 0.8, 3.9, 11.7, 0.8, "① 在庫あり=直す（小） / ② 在庫なし=戻す・畳む（中） / ③ 新規=増やす（大）",
    size=18, color=RGBColor(0xD9, 0xE2, 0xEC))
box(s, 0.8, 6.6, 11.7, 0.5, f"iMak Trading Japan  /  {datetime.date.today():%Y-%m-%d}", size=12, color=C_GREY)

# ── Slide 2: 3状態 × アクション ──
s = slide(); header(s, "3状態それぞれの効果的アクション")
column(s, 0.4, C_S1, "① 出品済・在庫あり\n＝直す", "2,605件", "小〜中",
       "• NO_CONVERT 高額\n   → 価格点検（値下げ/撤退）\n"
       "• NO_CLICK 高impr\n   → 画像/サムネ改善\n"
       "• NO_SEARCH（21日↑）\n   → タイトル front-load\n\n"
       "※ 低需要は何しても売れない\n   = 深追いしない（damage control）",
       "listing_funnel / amazon_v8_check / title_keyword_proposal")
column(s, 4.65, C_S2, "② 出品済・在庫なし\n＝戻す or 畳む", "1,992件", "中（実銭回収）",
       "• RESTOCK 再仕入れ\n   PSA13枚・G-SHOCK19件は黒字確定\n"
       "• 戻す前に ①タイトル ③価格 を準備\n"
       "• CULL 1,755\n   → 出品終了（アカウント健全化）\n\n"
       "※ 畳むは増収でなく衛生",
       "mercari_psa_resource / mercari_gshock_resource")
column(s, 8.9, C_S3, "③ 新規出品\n＝増やす（本丸）", "—", "大",
       "• 需要∩仕入れ可能∩自動化 の3つ\n   = UNIQLO-T / G-SHOCK定番 / Montbell\n"
       "• 既存パイプラインで面で量産\n"
       "• PORTER / PSA One Piece は\n   需要TOPでも拡大しない\n   （1点もの=タイムリー供給不可）\n\n"
       "※ 増収の本丸。出せる範囲に絞る",
       "tshirt_listing / gshock_to_csv / montbell_listing / demand_winners")

# ── Slide 3: 一言まとめ ──
s = slide(); header(s, "まとめ：どこに力を入れるか")
rows_txt = [
    (C_S1, "① 在庫あり（直す）", "深追いしない。高額の NO_CONVERT/NO_CLICK だけ手当て。効果は小。"),
    (C_S2, "② 在庫なし（戻す/畳む）", "再仕入れで実銭回収（PSA13/G-SHOCK19）＋ 死蔵1,755を畳んで健全化。効果は中。"),
    (C_S3, "③ 新規（増やす）★本丸", "『出せる定番』UNIQLO-T/G-SHOCK/Montbell を面で増やす。効果は大。"),
]
y = 1.5
for col, label, body in rows_txt:
    box(s, 0.7, y, 3.6, 1.1, label, size=16, bold=True, color=C_WHITE, fill=col,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(s, 4.5, y, 8.2, 1.1, body, size=14, color=C_DARK, fill=C_WHITE, line=col, anchor=MSO_ANCHOR.MIDDLE)
    y += 1.35
box(s, 0.7, 5.7, 12.0, 1.3,
    "結論: 最も効くのは ③（新規で出せる定番を面で増やす）。\n"
    "①は守り・②は回収/掃除。需要があっても『タイムリーに出せる』ものだけが増収につながる。",
    size=15, bold=True, color=C_PRIMARY, fill=RGBColor(0xE8, 0xEF, 0xD8), anchor=MSO_ANCHOR.MIDDLE)

out = os.path.join(DESK if os.path.isdir(DESK) else ".", f"出品3状態_効果的アクション_{datetime.date.today():%Y%m%d}.pptx")
prs.save(out)
print("出力:", out, "/ slides:", len(prs.slides._sldIdLst))
