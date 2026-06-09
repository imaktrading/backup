#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Item Specifics (itemsp) ロジック説明書 + itemsp CSV監査 → PPTX (デスクトップ)。
タイトル版 (make_title_logic_pptx.py) と同構成: ① 生成ロジック / ② CSV監査。"""
import datetime
import os
import sys
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DESK = r"C:\Users\imax2\OneDrive\デスクトップ"
FONT = "Yu Gothic UI"
C_PRIMARY = RGBColor(0x1F, 0x4E, 0x79)
C_ACCENT = RGBColor(0xC0, 0x50, 0x00)
C_GREEN = RGBColor(0x2E, 0x7D, 0x32)
C_DARK = RGBColor(0x20, 0x28, 0x30)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GREY = RGBColor(0x80, 0x80, 0x80)
C_BG = RGBColor(0xF4, 0xF6, 0xF9)
C_CODE = RGBColor(0x0B, 0x3D, 0x2E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = C_BG; bg.line.fill.background()
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def box(s, x, y, w, h, text, size=18, bold=False, color=C_DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    sp = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return sp


def code(s, x, y, w, h, text, size=14):
    sp = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor(0xEA, 0xF3, 0xEE); sp.line.color.rgb = C_GREEN
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15)
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = ln
        r.font.name = "Consolas"; r.font.size = Pt(size); r.font.color.rgb = C_CODE
    return sp


def header(s, title, sub=None):
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_PRIMARY; bar.line.fill.background()
    box(s, 0.5, 0.08, 12.3, 0.78, title, size=23, bold=True, color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        box(s, 0.5, 1.02, 12.3, 0.4, sub, size=13, color=C_ACCENT, bold=True)


def table(s, x, y, w, rows, col_w, fs=12, hfs=12, row_h=0.42):
    gt = s.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w),
                            Inches(row_h * len(rows))).table
    gt.first_row = False; gt.horz_banding = False
    for j, cw in enumerate(col_w):
        gt.columns[j].width = Inches(cw)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.margin_left = Inches(0.07); c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]; c.text_frame.word_wrap = True
            r = p.add_run(); r.text = str(val); r.font.name = FONT
            if i == 0:
                c.fill.solid(); c.fill.fore_color.rgb = C_PRIMARY
                r.font.size = Pt(hfs); r.font.bold = True; r.font.color.rgb = C_WHITE
            else:
                c.fill.solid(); c.fill.fore_color.rgb = C_WHITE if i % 2 else RGBColor(0xE8, 0xEE, 0xF5)
                r.font.size = Pt(fs); r.font.color.rgb = C_DARK
    return gt


# ============ 表紙 ============
s = slide()
hero = s.shapes.add_shape(1, 0, Inches(2.2), prs.slide_width, Inches(2.9))
hero.fill.solid(); hero.fill.fore_color.rgb = C_PRIMARY; hero.line.fill.background()
box(s, 0.8, 2.5, 11.7, 1.1, "Item Specifics ロジック 説明書", size=34, bold=True, color=C_WHITE)
box(s, 0.8, 3.7, 11.7, 0.9, "① Item Specifics の作り方 (catalog=SSOT)  /  ② Item Specifics のCSV監査 (CSV監査くん)",
    size=16, color=RGBColor(0xD9, 0xE2, 0xEC))
box(s, 0.8, 6.6, 11.7, 0.5, f"iMak Trading Japan  /  {datetime.date.today():%Y-%m-%d}  /  ※タイトル版と対の資料",
    size=12, color=C_GREY)

# ============ Part① divider ============
s = slide()
band = s.shapes.add_shape(1, 0, Inches(3.0), prs.slide_width, Inches(1.5))
band.fill.solid(); band.fill.fore_color.rgb = C_ACCENT; band.line.fill.background()
box(s, 0.8, 3.25, 11.7, 1.0, "Part ①  Item Specifics 生成ロジック", size=30, bold=True, color=C_WHITE)

# ============ 全体像: catalog SSOT + 3-case ============
s = slide(); header(s, "全体像 — catalog(公式DB)=SSOT を ID完全一致で引く", "タイトルと違い『加工しない』。公式値をそのまま/正規化/空欄の3択")
code(s, 0.5, 1.15, 12.3, 0.7, "catalog を ID 完全一致 lookup → 各 aspect を下表の3ケースで埋める (推測で埋めない)", size=14)
table(s, 0.5, 2.05, 12.3, [
    ["ケース", "eBayフィルタ種別", "やること", "例"],
    ["A", "SELECTION_ONLY", "catalog値を eBay公式値に正規化(whitelist)。無ければ空欄", "Rarity, Brand, Movement"],
    ["B", "FREE_TEXT", "catalog値をそのまま記入 (加工しない)", "Card Name, Set, Model"],
    ["C", "確証なし / 不明", "空欄。Country of Origin だけ 'Does not apply' を明示", "読めない型番=空欄"],
], [1.0, 3.0, 5.6, 2.7], row_h=0.66, fs=12)
box(s, 0.5, 4.95, 12.3, 1.7,
    "■ 鉄則 (出品の正確性原則)\n"
    "・ID 完全一致 lookup のみ。名前検索フォールバック禁止 → ID不一致は reject\n"
    "・確証なき情報は空欄 (推測で埋めない)。空欄に放置すると eBay AI が勝手に補完するので Country は明示的に 'Does not apply'\n"
    "・加工 (SEO アレンジ) は title / description に閉じ込め、Item Specifics は公式値のまま",
    size=13, color=C_DARK)

# ============ whitelist 自己修正ループ ============
s = slide(); header(s, "① SELECTION_ONLY の正規化 — whitelist_registry 自己修正ループ")
box(s, 0.5, 1.15, 12.3, 0.9,
    "eBay の SELECTION_ONLY フィールドは『公式ドロップダウンの正規値』と完全一致しないとフィルタに載らない。\n"
    "whitelist_registry (カテゴリ別の正規値+誤表記マップ) で Claude 出力を検証・正規化する。", size=14, bold=True, color=C_DARK)
code(s, 0.5, 2.2, 12.3, 1.0,
     "normalized, violations = validate_and_normalize(item_specifics, category)\n"
     "if violations:  feedback = build_retry_feedback(violations) → Claude に再リクエスト", size=13)
table(s, 0.5, 3.4, 12.3, [
    ["要素", "意味"],
    ["values", "有効値リスト (完全一致)。strict=True ならリスト外は違反"],
    ["normalize", "誤表記→正規値 の自動修正マップ (例 'UNIQLO UT'→'Uniqlo')"],
    ["multi", "カンマ区切り複数値を許可するか"],
], [2.2, 9.6], row_h=0.5, fs=12)
box(s, 0.5, 5.3, 12.3, 1.0,
    "例: tshirt Brand は values=['Uniqlo'] / normalize={'UNIQLO UT':'Uniqlo'}。\n"
    "→ 'UNIQLO UT' は eBayに無いブランド名なので 'Uniqlo' に矯正 (フィルタ不ヒット=売上機会喪失を防ぐ)。",
    size=13, color=C_ACCENT, bold=True)

# ============ Vision gap-fill + fail-closed ============
s = slide(); header(s, "Vision は穴埋め専用 / fail-closed 設計")
box(s, 0.5, 1.2, 12.3, 2.2,
    "■ Vision (画像AI) = gap-fill 専用\n"
    "・catalog に在る値は Vision で上書きしない。catalog が空の項目だけ Vision で補完。\n"
    "・型番(MPN)は画像のタグから読めた時のみ。公式サイトからの推定は不可 (同名で複数型番あり)。\n\n"
    "■ fail-closed (Precision 100% / Recall は諦める)\n"
    "・認識できない → 出品しない (skip)。認識を間違えて出品 → 絶対NG (SNAD/Defect Rate)。\n"
    "・確証なき情報は空欄。「網羅性が低い/出品数が減る」は受け入れる。", size=14, color=C_DARK)
box(s, 0.5, 5.7, 12.3, 0.9,
    "→ Item Specifics は『公式 catalog の事実』だけで構成。出品の都度カタログに依頼を減らすため、catalog を充実させるのが本筋。",
    size=13, color=C_ACCENT, bold=True)

# ============ カテゴリ別 主要 Item Specifics ============
s = slide(); header(s, "カテゴリ別 — 主要 Item Specifics (catalog 由来)")
table(s, 0.5, 1.2, 12.3, [
    ["カテゴリ", "主要フィールド (例)"],
    ["PSA TCG", "Game / Card Name / Character / Set / Card Number / Rarity / Card Type / Finish / Country"],
    ["G-SHOCK", "Brand=Casio / Department=Men / Style / Display / Case&Band Color / Material / Movement / Water Resistance / Model(シリーズ)"],
    ["一番くじ", "Character / Theme=Anime & Manga / Type=Figure / Franchise / Brand=Bandai / Country"],
    ["Porter", "Brand=Porter / Department / Style / Color / Material / Size"],
    ["Tomica", "Brand=Tomica / Vehicle Type / Make / Material=Diecast / Scale(1:64) / Color / Country=Japan"],
    ["UNIQLO", "Brand=Uniqlo / Product Line=Uniqlo UT / Theme / Character Family / Size / Color / Department"],
    ["リール", "Brand / Model / Gear Ratio / Item Weight / Maximum Drag / Reel Type"],
    ["Workman", "Brand / Type / Material / Style / Activity / Season / Department"],
], [1.7, 10.6], row_h=0.55, fs=11.5)
box(s, 0.5, 6.5, 12.3, 0.4, "値は全て catalog 公式由来 (eBay正規化 or そのまま or 空欄)。Country of Origin は不明時 'Does not apply'。",
    size=12, color=C_GREY)

# ============ Part② divider ============
s = slide()
band = s.shapes.add_shape(1, 0, Inches(3.0), prs.slide_width, Inches(1.5))
band.fill.solid(); band.fill.fore_color.rgb = C_GREEN; band.line.fill.background()
box(s, 0.8, 3.25, 11.7, 1.0, "Part ②  Item Specifics のCSV監査 (CSV監査くん)", size=28, bold=True, color=C_WHITE)

# ============ 監査の考え方 ============
s = slide(); header(s, "監査の考え方 — 取得済 eBay公式フィルタJSON を武器にする", "持っている武器(公式Aspects)を活かす。offline・API不要")
box(s, 0.5, 1.6, 12.3, 1.0,
    "eBay の getItemAspectsForCategory で取得した公式 Aspects JSON (各 aspect の values / 必須・推奨 / SELECTION_ONLY|FREE_TEXT)\n"
    "と CSV の Item Specifics を突き合わせる。生成と同じ公式値を基準にするので二重基準にならない。", size=14, bold=True, color=C_DARK)
table(s, 0.5, 2.9, 12.3, [
    ["検査", "見るもの", "処置"],
    ["① 正規値ゲート", "SELECTION_ONLY の値が公式許容リスト外 → フィルタ不ヒット", "報告 (SEO)"],
    ["② 未充足ゲート", "必須/推奨 aspect が『列無し or 全行空』 → 検索性の取りこぼし", "報告 (SEO)"],
    ["③ 必須spec空", "validate_row が必須Item Specific空を検出", "除外+カタログ依頼"],
    ["④ ドリフト監視", "手動whitelist vs 公式フィルタ の値ズレ (週次バッチ)", "whitelist修正"],
], [2.4, 7.0, 2.9], row_h=0.58, fs=12)
box(s, 0.5, 6.4, 12.3, 0.5, "特殊値 ('Does not apply'/'N/A' 等 eBay普遍の opt-out) は許容＝誤検出しない。",
    size=12, color=C_GREY)

# ============ ① 正規値ゲート 詳細 ============
s = slide(); header(s, "① 正規値ゲート — SELECTION_ONLY 許容外の検出")
code(s, 0.5, 1.2, 12.3, 0.95,
     "公式 aspect_mode==SELECTION_ONLY の各フィールドで、CSVの値が values[] に無ければ:\n"
     "  「'Rarity'='Ultra Rare' が公式フィルタ許容値外(SELECTION_ONLY)→フィルタ不ヒット」", size=13)
box(s, 0.5, 2.4, 12.3, 1.6,
    "・eBayの『絞り込みドロップダウン』はこの正規値しか拾わない → 表記揺れ/自由文字は検索に載らない\n"
    "・例: Tシャツ Theme='Anime & Manga' は無効、正解は 'Anime' / Brand='UNIQLO UT' は無効、正解は 'Uniqlo'\n"
    "・行ごとに検査 (上限30件で打切り、報告のみ=CSVは触らない)\n"
    "・eBay普遍の特殊値 (Does not apply / N/A / Unbranded 等) は許容リストに無くても誤検出しない", size=13, color=C_DARK)
box(s, 0.5, 4.2, 12.3, 0.9,
    "→ 生成側 (whitelist_registry) と監査側 (公式JSON) が同じ正規値を基準にするので、\n"
    "  生成で正規化漏れがあれば監査が拾い、whitelist を公式値へ寄せる (ドリフト監視④と連動)。",
    size=13, color=C_ACCENT, bold=True)

# ============ ②③④ 詳細 ============
s = slide(); header(s, "② 未充足(SEO機会) / ③ 必須空(除外+依頼) / ④ ドリフト監視")
box(s, 0.5, 1.15, 12.3, 1.5,
    "■ ② 未充足ゲート (SEO機会)\n"
    "・eBay が必須/推奨としている aspect が『列が無い』or『全行空』→ 埋めれば検索性UP の機会として報告\n"
    "・値は足さない (catalog に無ければ catalog 拡充が筋)。報告のみ。", size=13, color=C_DARK)
box(s, 0.5, 2.8, 12.3, 1.5,
    "■ ③ 必須Item Specific 空 → 除外 + カタログ依頼\n"
    "・必須項目が空の行は fail-closed で除外 (誤出品しない)。catalog が空なら Catalog/Harvest に修正依頼。\n"
    "・CSV に値を書き足さない (捏造禁止)。catalog=SSOT を維持。", size=13, color=C_DARK)
box(s, 0.5, 4.45, 12.3, 1.5,
    "■ ④ ドリフト監視 (週次バッチ whitelist_official_drift)\n"
    "・手動 whitelist_registry の値が、取得済の公式 SELECTION_ONLY 許容リストから外れてないか照合。\n"
    "・ズレ (持ち腐れ) を検出 → whitelist を公式値へ修正。毎週月曜の Catalog Integrity バッチで自動実行。", size=13, color=C_DARK)

# ============ まとめ ============
s = slide(); header(s, "まとめ")
box(s, 0.6, 1.4, 12.1, 5.0,
    "① Item Specifics 生成ロジック\n"
    "   ・catalog(公式DB)=SSOT を ID完全一致で引く。加工しない (SEOはtitle/descに閉じる)\n"
    "   ・3ケース: SELECTION_ONLY=whitelist正規化 / FREE_TEXT=そのまま / 確証なし=空欄(Country='Does not apply')\n"
    "   ・Vision=穴埋め専用 (catalog上書き禁止) / fail-closed (ID不一致=reject、推測禁止)\n\n"
    "② Item Specifics のCSV監査 (CSV監査くん) — 取得済 公式Aspects JSON が武器\n"
    "   ・①正規値ゲート(SELECTION_ONLY許容外) ②未充足(必須/推奨が空=SEO機会) ③必須空(除外+カタログ依頼) ④ドリフト監視(週次)\n"
    "   ・値の捏造はしない。違反は報告 or 除外+依頼。生成と監査が同じ公式値基準=二重基準にならない\n\n"
    "→ 本筋: 出品の都度カタログに依頼するのを減らすため、catalog を充実させ生成が公式値で満たすこと",
    size=14, color=C_DARK)

out = os.path.join(DESK, f"ItemSpecificsロジック説明書_{datetime.date.today():%Y%m%d}.pptx")
prs.save(out)
print(f"✅ 保存: {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
